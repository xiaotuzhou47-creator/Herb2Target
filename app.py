# -*- coding: utf-8 -*-
import sqlite3, json, math, uuid, os, time, urllib.request, urllib.parse, hashlib
import xml.etree.ElementTree as ET
from flask import Flask, jsonify, request, Response
from flask_cors import CORS

app = Flask(__name__); CORS(app)

# 对接路由必须最早注册
@app.route("/api/v1/fix_pocket", methods=["POST"])
def fix_pocket():
    data = request.get_json() or {}
    # Admin endpoint — intended for database maintenance, not exposed in the web UI
    uniprot = data.get("uniprot","")
    score = data.get("pocket_score")
    has = data.get("has_pocket")
    if not uniprot: return jsonify({"error":"need uniprot"}),400
    conn = get_db(); c = conn.cursor()
    c.execute("UPDATE targets SET pocket_score=?, has_pocket=? WHERE uniprot_id=?",(score,has,uniprot))
    conn.commit(); conn.close()
    return jsonify({"status":"ok","updated":c.rowcount})


@app.route("/api/v1/dock", methods=["POST"])
def dock_molecule():
    # v2: RDKit+largeMol+1CX2 + force_blind对比
    data = request.get_json() or {}
    s = data.get("smiles","").strip(); u = data.get("uniprot","").strip()
    fb = data.get("force_blind", False)
    bs = data.get("box_size", 28)
    rec = data.get("receptor_pdb","").strip()
    if not s or not u: return jsonify({"error":"Need smiles+uniprot"}),400
    try:
        import sys
        for k in list(sys.modules.keys()):
            if 'docking' in k: del sys.modules[k]
        from docking import dock_fast
        return jsonify(dock_fast(s,u,receptor_pdb=rec if rec else None, force_blind=fb, box_size=bs))
    except Exception as e:
        return jsonify({"status":"error","message":str(e)})

@app.route("/api/v1/dock_all", methods=["POST"])
def dock_all_targets():
    data = request.get_json() or {}
    smiles = data.get("smiles","").strip()
    if not smiles: return jsonify({"error":"Need SMILES"}),400
    pred = predict_molecule_targets(smiles)
    all_targets = pred.get("targets",[])
    # 优先PS>=0.5的靶点，但如果没有足够的，对has_pocket=True但PS<0.5的做盲对接
    pocket_targets = [t for t in all_targets if t.get("has_pocket") and (t.get("pocket_score") or 0) >= 0.5][:5]
    if len(pocket_targets) < 3:
        # 补充has_pocket但PS低的靶点（盲对接）
        weak_targets = [t for t in all_targets if t.get("has_pocket") and (t.get("pocket_score") or 0) < 0.5][:3]
        pocket_targets = (pocket_targets + weak_targets)[:5]
    targets = pocket_targets
    from docking import dock_fast
    results = []
    for t in targets:
        r = dock_fast(smiles, t["uniprot_id"])
        results.append({"gene":t["gene_name"],"uniprot":t["uniprot_id"],"affinity":r.get("affinity"),"status":r.get("status"),"method":r.get("method",""),"pose_pdbqt":r.get("pose_pdbqt","")})
    return jsonify({"smiles":smiles,"results":results})
BASE = os.path.dirname(os.path.realpath(__file__))
DB_PATH = os.path.join(BASE, "tcm_demo.db")

# === SMILES Parser & ECFP ===
ATOM_NUM = {'C':6,'N':7,'O':8,'S':16,'P':15,'F':9,'Cl':17,'Br':35,'I':53,'B':5,'Si':14,'Se':34,'c':6,'n':7,'o':8,'s':16}
BOND_ORDER = {'-':1,'=':2,'#':3,':':1.5}

def _tokenize(s):
    t=[];i=0
    while i<len(s):
        c=s[i]
        if c in '.-=#:()[]' or c.isdigit():
            if c=='%' and i+2<len(s) and s[i+1:i+3].isdigit():t.append(s[i:i+3]);i+=3;continue
            t.append(c);i+=1
        elif c.isalpha():
            if i+1<len(s) and s[i+1].islower():t.append(s[i:i+2]);i+=2
            else:t.append(c);i+=1
        else:t.append(c);i+=1
    return t

def _parse_smiles(s):
    tokens=_tokenize(s);atoms=[];bonds=[];stack=[];rc={};ca=None;cb='-'
    def add_atom(e,ar=False):
        nonlocal ca,cb
        atoms.append({'el':e,'ch':0,'ar':ar,'h':0,'idx':len(atoms)})
        if ca is not None:bonds.append({'s':ca,'d':len(atoms)-1,'o':BOND_ORDER.get(cb,1)})
        ca=len(atoms)-1;cb='-'
    i=0
    while i<len(tokens):
        t=tokens[i]
        if t=='(':
            stack.append(ca);cb='-'
        elif t==')':
            if stack:ca=stack.pop()
            cb='-'
        elif t in '.-=#:':
            cb=t if t!='.' else '-'
            if t=='.':ca=None
        elif t.isdigit() or (t.startswith('%') and t[1:].isdigit()):
            rn=int(t.replace('%',''))
            if rn in rc:
                p=rc.pop(rn)
                if p!=ca:bonds.append({'s':p,'d':ca,'o':BOND_ORDER.get(cb,1)})
                cb='-'
            else:
                if ca is not None:rc[rn]=ca
                cb='-'
        elif t=='[':
            bkt=[];i+=1;d=1
            while i<len(tokens) and d>0:
                if tokens[i]=='[':d+=1
                elif tokens[i]==']':d-=1
                if d>0:bkt.append(tokens[i])
                i+=1
            e='C';ch=0;ar=False
            for bt in bkt:
                if bt in ATOM_NUM:e=bt;ar=bt.islower() and bt in 'cno'
                elif bt=='+':ch+=1
                elif bt=='-':ch-=1
            atoms.append({'el':e,'ch':ch,'ar':ar,'h':0,'idx':len(atoms)})
            if ca is not None:bonds.append({'s':ca,'d':len(atoms)-1,'o':BOND_ORDER.get(cb,1)})
            ca=len(atoms)-1;cb='-'
        else:add_atom(t,t.islower() and t in 'cno')
        i+=1
    return atoms,bonds

def _classify_smiles(smiles):
    """化学类型预分类"""
    c = smiles.count('C') + smiles.count('c')
    o = smiles.count('O') + smiles.count('o')
    n = smiles.count('N') + smiles.count('n')
    if '[N+]' in smiles or 'n+]' in smiles: return 'alkaloid'
    if n >= 1 and ('CN1CCC' in smiles or 'CN1CC' in smiles): return 'alkaloid'
    flav_sigs = ['C1=CC(=O)C2=C','C1C(=O)C2=C','C1=CC(=O)c2c','C=CC(=O)c1c(O)','OC=CC(=O)c1','C(=O)C=COc1']
    if any(s in smiles for s in flav_sigs) and o >= 3: return 'flavonoid'
    if 'c1coc2c' in smiles and o >= 4 and 'C(=O)' in smiles: return 'flavonoid'
    if 'C(=O)c1c' in smiles and 'C(=O)c2c' in smiles and c >= 14: return 'anthraquinone'
    if c >= 15 and o <= 4 and n == 0 and any(s in smiles for s in ['CC(C)C','CC1(C)CC','C=C(C)C']): return 'terpenoid'
    if o >= 2 and c <= 12 and n == 0: return 'phenolic'
    # polyphenol covers flavonoids, stilbenoids, diarylheptanoids
    if o >= 2 and c >= 14: return 'polyphenol'
    if n >= 1 and 'c1c' in smiles: return 'alkaloid'
    return 'other'

def ecfp_fingerprint(smiles,n_bits=2048,mr=2):
    clean=smiles.replace('@@','').replace('@','').replace('/','').replace('\\','')
    try:atoms,bonds=_parse_smiles(clean)
    except:return None
    if not atoms:return None
    adj=[[] for _ in range(len(atoms))]
    for b in bonds:adj[b['s']].append((b['d'],b['o']));adj[b['d']].append((b['s'],b['o']))
    for i,a in enumerate(atoms):a['deg']=len(adj[i])
    inv=[];bits=[0]*n_bits
    for a in atoms:an=ATOM_NUM.get(a['el'],6);inv.append(['%d_%d_%d'%(an,a['deg'],a['ch'])])
    for r in range(1,mr+1):
        for i in range(len(atoms)):
            nb=[];e=''
            for ni,bo in adj[i]:
                if len(inv[ni])>=r:nb.append('%s_%d'%(inv[ni][r-1],int(bo)))
            nb.sort();e=inv[i][r-1]+'|'+''.join(nb)
            if len(inv[i])<=r:inv[i].append(e)
            else:inv[i][r]=e
    for fl in inv:
        for f in fl:
            h=int(hashlib.sha256(f.encode()).hexdigest(),16)%n_bits;bits[h]=1
    return bits

def tanimoto_similarity(a,b):
    if not a or not b:return 0.0
    n=min(len(a),len(b));ab=sum(1 for i in range(n) if a[i] and b[i])
    ob=sum(1 for i in range(n) if a[i] or b[i])
    return ab/max(ob,1)

_FC={}
def build_fingerprint_cache(force=False):
    global _FC
    if _FC and not force:return
    # 仅加载IT_MAP中存在的化合物的指纹（避免无关参考化合物污染）
    conn=sqlite3.connect(DB_PATH);conn.row_factory=sqlite3.Row;c=conn.cursor()
    c.execute("SELECT id,name,smiles FROM ingredients WHERE smiles IS NOT NULL AND smiles!=''")
    rows=c.fetchall();conn.close()
    valid_names = set(IT_MAP.keys())
    cnt=0
    for rid,name,smiles in rows:
        if name in _FC:continue
        if name not in valid_names: continue  # 只加载IT_MAP中的化合物
        fp=ecfp_fingerprint(smiles)
        if fp:_FC[name]=fp;cnt+=1
    print('  ECFP cache: %d compounds (IT_MAP-filtered)'%cnt)

_BG_STATS = None

def _compute_background():
    global _BG_STATS
    if _BG_STATS: return
    if not _FC: build_fingerprint_cache()
    fps=[(n,fp) for n,fp in _FC.items() if fp is not None]
    if len(fps)<10: _BG_STATS=(0.15,0.1);return
    import random
    scores=[]
    for _ in range(5000):
        a=fps[random.randint(0,len(fps)-1)][1]
        b=fps[random.randint(0,len(fps)-1)][1]
        scores.append(tanimoto_similarity(a,b))
    mu=sum(scores)/len(scores)
    sigma=(sum((s-mu)**2 for s in scores)/len(scores))**0.5
    _BG_STATS=(mu,sigma)
    print(f'  Background: mu={mu:.3f} sigma={sigma:.3f}')

def predict_molecule_targets(smiles,top_n=5):
    # SMILES 安全过滤
    if not isinstance(smiles,str) or len(smiles)>2000 or len(smiles)<3:
        return{'error':'Invalid SMILES','targets':[]}
    # 过滤非 ASCII 字符
    clean_smi = ''.join(c for c in smiles if ord(c) < 128)
    if not clean_smi: return{'error':'Only ASCII SMILES allowed','targets':[]}
    q_fp=ecfp_fingerprint(clean_smi)
    if q_fp is None:return{'error':'SMILES format error','targets':[]}
    # 每次请求检查IT_MAP是否需要重新加载
    import os as _os
    _itmap_mtime = _os.path.getmtime(_os.path.join(BASE, 'itmap_raw.json'))
    if not hasattr(build_fingerprint_cache, '_itmap_mtime') or build_fingerprint_cache._itmap_mtime < _itmap_mtime:
        build_fingerprint_cache._itmap_mtime = _itmap_mtime
        # 强制重建_FC缓存
        build_fingerprint_cache(force=True)
    
    if not _FC:build_fingerprint_cache()
    
    # 化学类型预分类：只在匹配的子库中搜索
    ctype = _classify_smiles(clean_smi)
    # 加载分层IT_MAP
    if not hasattr(build_fingerprint_cache, '_stratified_itmap'):
        with open(os.path.join(BASE, 'itmap_stratified.json'), 'r') as _sf:
            build_fingerprint_cache._stratified_itmap = json.load(_sf)
    sub_lib = build_fingerprint_cache._stratified_itmap.get(ctype, {})
    sub_lib_names = set(sub_lib.keys())
    
    # 不使用分层，全库搜索。频率归一化用更强的log10
    sims=[]
    for n,fp in _FC.items():
        if fp is None:continue
        s=tanimoto_similarity(q_fp,fp)
        if s > 0.99: continue  # 排除自我匹配
        sims.append((n,s))
    sims.sort(key=lambda x:x[1],reverse=True)
    ts=sims[:top_n];best=ts[0][1] if ts else 0;lc=best<0.4
    if not lc:ts=[(n,s) for n,s in ts if s>0.25]
    # Top-k加权投票：取前k个最相似参考化合物，相似度加权投票
    # 然后除以靶点频率（TF-IDF风格归一化），消除高频靶点优势
    K = min(len(ts), 30)  # top-30参考化合物
    ts_k = ts[:K]
    
    # 计算靶点频率（IDF）
    bg_freq={};total_compounds=len(IT_MAP)
    for gs in IT_MAP.values():
        for g in gs:bg_freq[g]=bg_freq.get(g,0)+1
    
    # 加权投票：每个参考化合物按相似度投票给其关联靶点
    conn=sqlite3.connect(DB_PATH);conn.row_factory=sqlite3.Row;c=conn.cursor()
    c.execute('SELECT id,gene_name FROM targets');at={r['gene_name']:r['id'] for r in c.fetchall()};conn.close()
    sc={};sd=[]
    for n,s in ts_k:
        gs=IT_MAP.get(n,[])
        sd.append({'name':n,'similarity':round(s,4),'target_count':len(gs)})
        weight = s * s  # 平方加权：高相似度参考化合物权重更大
        for g in gs:
            if at.get(g):
                sc[g]=sc.get(g,0)+weight
    
    # Log频率归一化：除以log2(1+freq)消除高频靶点天然优势
    for g in list(sc.keys()):
        freq = bg_freq.get(g, 1)
        sc[g] = sc[g] / (1.0 + math.log2(max(freq, 1)))
    
    # 加载口袋可药性评分(PS)和对接基准AUC
    conn=sqlite3.connect(DB_PATH);conn.row_factory=sqlite3.Row;c=conn.cursor()
    c.execute('SELECT gene_name,pocket_score FROM targets WHERE pocket_score IS NOT NULL')
    ps_scores = {r['gene_name']:r['pocket_score'] for r in c.fetchall() if r['pocket_score'] and r['pocket_score'] >= 0}
    c.execute('SELECT gene_name,auc FROM docking_benchmark')
    docking_auc = {r['gene_name']:r['auc'] for r in c.fetchall()}
    conn.close()
    
    # PS调整因子：PS>=0.5的靶点获得排名提权，PS<0.5的降权
    # PS权重因子：10对验证集上1.0达到最优Top-5=90%
    ps_weight = 1.0
    for g in list(sc.keys()):
        ps = ps_scores.get(g, None)
        if ps is not None:
            # PS调整：(PS-0.5)为正则提权，为负则降权
            adj = 1.0 + ps_weight * (ps - 0.5)
            sc[g] = sc[g] * max(adj, 0.3)  # 不低于0.3，避免PS=0的靶点被完全清零
    
    st=sorted(sc.items(),key=lambda x:x[1],reverse=True)
    nt=len(ts_k) if ts_k else 1;rt=[]
    for g,ss in st:
        conn=sqlite3.connect(DB_PATH);conn.row_factory=sqlite3.Row;c2=conn.cursor()
        c2.execute('SELECT * FROM targets WHERE gene_name=?',(g,));row=c2.fetchone();conn.close()
        if row:
            ev=round(ss/nt,4);nv=row['novelty_score']
            combined=round(ev*nv,4)
            auc_val = docking_auc.get(row['gene_name'], None)
            rt.append({'gene_name':row['gene_name'],'uniprot_id':row['uniprot_id'],
                'protein_name':row['protein_name'],'pubmed_total':row['pubmed_total'],
                'pubmed_disease':row['pubmed_disease'],'novelty_score':nv,
                'has_pocket':bool(row['has_pocket'] and (row['pocket_score'] or 0) >= 0),'pocket_score':row['pocket_score'] if (row['pocket_score'] or 0) >= 0 else None,
                'docking_auc':auc_val,
                'confidence':ev,'combined_score':combined,
                'target_class':TARGET_CLASSES.get(row['gene_name'],'其他')})
    rt.sort(key=lambda x:x['combined_score'],reverse=True)
    return{'similar_compounds':sd,'target_count':len(rt),'targets':rt,'ppi_inferred':[]}

# PPI-based novel target inference (guilt-by-association)
def infer_ppi_targets(target_scores, top_similar):
    """PPI网络推断：STRING置信度加权，guilt-by-association"""
    ppi_undirected={}
    for g1,partners in PPI_MAP.items():
        for g2 in partners:
            ppi_undirected.setdefault(g1,set()).add(g2)
            ppi_undirected.setdefault(g2,set()).add(g1)
    inferred={}
    conn=sqlite3.connect(DB_PATH);conn.row_factory=sqlite3.Row;c=conn.cursor()
    c.execute('SELECT id,gene_name,novelty_score FROM targets WHERE novelty_score>=0.6')
    novel_tgts={r['gene_name']:r['id'] for r in c.fetchall()}
    direct_genes=set(target_scores.keys())
    for g1 in direct_genes:
        if g1 not in ppi_undirected:continue
        evidence=target_scores.get(g1,0)
        for g2 in ppi_undirected[g1]:
            if g2 not in novel_tgts:continue
            if g2 in direct_genes:continue
            if g2 not in inferred:inferred[g2]={'score':0,'total_sim':0,'conn':[]}
            c.execute('SELECT score FROM ppi WHERE(t1 IN(SELECT id FROM targets WHERE gene_name=?)AND t2 IN(SELECT id FROM targets WHERE gene_name=?))OR(t2 IN(SELECT id FROM targets WHERE gene_name=?)AND t1 IN(SELECT id FROM targets WHERE gene_name=?))',(g1,g2,g1,g2))
            ppi_row=c.fetchone()
            ppi_conf=(ppi_row['score']/1000.0)if ppi_row else 0.7
            inferred[g2]['score']+=evidence*ppi_conf
            inferred[g2]['total_sim']+=ppi_conf
            inferred[g2]['conn'].append(g1)
    result=[]
    nt=len(top_similar)if top_similar else 1
    for gene,data in sorted(inferred.items(),key=lambda x:x[1]['score']/max(x[1]['total_sim'],0.001),reverse=True)[:5]:
        c.execute('SELECT*FROM targets WHERE gene_name=?',(gene,))
        row=c.fetchone()
        if row:
            avg_conf=data['score']/max(data['total_sim'],0.001)
            ev=round(avg_conf/nt,4);nv=row['novelty_score']
            combined=round(ev*nv,4)
            result.append({'gene_name':row['gene_name'],'uniprot_id':row['uniprot_id'],
                'protein_name':row['protein_name'],'pubmed_total':row['pubmed_total'],
                'pubmed_disease':row['pubmed_disease'],'novelty_score':nv,
                'has_pocket':bool(row['has_pocket'] and (row['pocket_score'] or 0) >= 0),'pocket_score':row['pocket_score'] if (row['pocket_score'] or 0) >= 0 else None,
                'ps_warning': ('Low-confidence AlphaFold structure; interpret with caution' if (row['pocket_score'] or 0) > 0 and (row['pocket_score'] or 0) < 0.5 else (None if (row['pocket_score'] or 0) >= 0 else 'Structure unavailable, pocket not assessed')),
                'confidence':ev,'combined_score':combined,
                'target_class':TARGET_CLASSES.get(row['gene_name'],'其他'),
                'source':'PPI推断','ppi_partners':','.join(data['conn'][:3])})
    conn.close()
    return result

# Auto-expand IT_MAP for new ingredients on startup
def auto_expand_itmap():
    """对新加入的成分自动生成靶点映射（自举）"""
    if not _FC: build_fingerprint_cache()
    added=0
    # 手写成分集合
    hand_names = set(r[1] for r in INGREDIENTS_DATA if r[0] != '参考库')
    for name in hand_names:
        if name in IT_MAP: continue
        if name not in _FC: continue
        q_fp=_FC[name]
        sims=[]
        for n,fp in _FC.items():
            if n==name or fp is None or n not in hand_names: continue
            s=tanimoto_similarity(q_fp,fp);sims.append((n,s))
        sims.sort(key=lambda x:x[1],reverse=True)
        ts=sims[:5]
        sc={}
        for n,s in ts:
            for g in IT_MAP.get(n,[]):sc[g]=sc.get(g,0)+s
        genes=[g for g,_ in sorted(sc.items(),key=lambda x:x[1],reverse=True)[:8]]
        if genes:
            IT_MAP[name]=genes;added+=1
            # 写入数据库
            try:
                conn=sqlite3.connect(DB_PATH)
                conn.row_factory=sqlite3.Row
                c=conn.cursor()
                c.execute('SELECT id FROM ingredients WHERE name=?',(name,))
                ing=c.fetchone()
                if ing:
                    for g in genes:
                        c.execute('SELECT id FROM targets WHERE gene_name=?',(g,))
                        tgt=c.fetchone()
                        if tgt:
                            c.execute('INSERT OR IGNORE INTO ingredient_targets(ingredient_id,target_id)VALUES(?,?)',(ing['id'],tgt['id']))
                conn.commit();conn.close()
            except: pass
    print('  Auto-expanded IT_MAP: %d new compounds' % added)
def serve_html():
    p=os.path.join(BASE,"templates","index.html")
    with open(p,"rb") as f:c=f.read()
    return Response(c,mimetype="text/html",headers={"Cache-Control":"no-cache,no-store,must-revalidate"})

@app.route("/pptx_generator")
def pptx_generator():
    p=os.path.join(BASE,"templates","pptx_simple.html")
    if os.path.exists(p):
        with open(p,"rb") as f:c=f.read()
        return Response(c,mimetype="text/html",headers={"Cache-Control":"no-cache,no-store,must-revalidate"})
    return Response("Not found",status=404)

@app.route("/pptx_full")
def pptx_full():
    p=os.path.join(BASE,"templates","pptx_generator.html")
    if os.path.exists(p):
        with open(p,"rb") as f:c=f.read()
        return Response(c,mimetype="text/html",headers={"Cache-Control":"no-cache,no-store,must-revalidate"})
    return Response("Not found",status=404)

@app.route("/api/gen_finance")
def gen_finance():
    import sys
    sys.path.insert(0, os.path.join(BASE, 'templates'))
    from FinanceGen import build
    desktop = os.path.join(os.path.expanduser('~'), 'Desktop')
    outpath = os.path.join(desktop, '大创', '财务预测表.docx')
    data = build()
    with open(outpath, 'wb') as f:
        f.write(data)
    return Response(f'OK: saved to {outpath}')

@app.route("/api/fix_plan_gen")
def fix_plan_gen():
    import re
    path = os.path.join(BASE, 'templates', 'ProjectPlanGen.py')
    with open(path, 'r', encoding='utf-8') as f:
        code = f.read()
    
    # 1. 标题编号：一、二、三... → 1 2 3...
    cn_nums = {'一':'1','二':'2','三':'3','四':'4','五':'5','六':'6','七':'7','八':'8','九':'9','十':'10'}
    for cn, num in cn_nums.items():
        code = code.replace(f"'{cn}、", f"'{num} ")
    
    # 2. 百分数：百分之X → X%
    pct_map = {
        '百分之七十':'70%','百分之九十二':'92%','百分之八十六':'86%','百分之七十八':'78%',
        '百分之三十':'30%','百分之八十五':'85%','百分之三十三':'33%','百分之六十六点七':'66.7%',
        '百分之五十':'50%','百分之十':'10%','百分之十三点七':'13.7%','百分之五':'5%',
        '百分之八十八':'88%','百分之九十':'90%','百分之八十二':'82%','百分之三十一':'31%',
        '百分之五十一':'51%'
    }
    for old, new in pct_map.items():
        code = code.replace(old, new)
    
    # 3. 零点X → 0.X
    code = code.replace('零点一八','0.18').replace('零点五','0.5')
    
    # 4. 第一/第二/第三/第四/第五/第六 → 1/2/3/4/5/6 （在段落开头位置）
    # 但要注意不要破坏标题里的数字
    
    with open(path, 'w', encoding='utf-8') as f:
        f.write(code)
    return Response(f'Fixed. File size: {len(code)} chars')

@app.route("/api/fix_plan_gen2")
def fix_plan_gen2():
    path = os.path.join(BASE, 'templates', 'ProjectPlanGen.py')
    with open(path, 'r', encoding='utf-8') as f:
        code = f.read()
    
    # Replace first/second/etc in text with 1/2/3
    replacements = [
        ("'第一是", "'1是"),
        ("'第二是", "'2是"),
        ("'第三是", "'3是"),
        ("'第四是", "'4是"),
        ("'第五是", "'5是"),
        ("第一是多", "1是多"),
        ("第二是苗药", "2是苗药"),
        ("第三是新药", "3是新药"),
        ("第四是数据", "4是数据"),
        ("第五是AI", "5是AI"),
        ("第一是多技术栈", "1是多技术栈"),
        ("第一是工具碎片化", "1是工具碎片化"),
        ("第二是流程耗时长", "2是流程耗时长"),
        ("第三是技术门槛高", "3是技术门槛高"),
        ("第四是结果不可复现", "4是结果不可复现"),
        ("第五是数据孤岛", "5是数据孤岛"),
        ("第六是民族药物空白", "6是民族药物空白"),
        # Also handle the section numbering within paragraphs
        ("第一是多技术栈融合能力", "1是多技术栈融合能力"),
        ("第二是数据整合壁垒", "2是数据整合壁垒"),
        ("第三是算法优化壁垒", "3是算法优化壁垒"),
        ("第四是苗药数据壁垒", "4是苗药数据壁垒"),
    ]
    for old, new in replacements:
        code = code.replace(old, new)
    
    with open(path, 'w', encoding='utf-8') as f:
        f.write(code)
    return Response(f'Fix2 done. Size: {len(code)}')

@app.route("/api/read_final")
def read_final():
    from docx import Document
    path = os.path.join(os.path.expanduser('~'), 'Desktop', '大创', '项目计划书_最终版.docx')
    if not os.path.exists(path):
        return Response('File not found', status=404)
    doc = Document(path)
    texts = []
    for p in doc.paragraphs:
        if p.text.strip():
            texts.append(p.text[:200])
    return Response('\n'.join(texts[:100]))

@app.route("/api/read_doc_structure")
def read_doc_structure():
    from docx import Document
    path = os.path.join(os.path.expanduser('~'), 'Desktop', '大创', '项目计划书_最终版.docx')
    if not os.path.exists(path):
        return Response('File not found', status=404)
    doc = Document(path)
    result = []
    for i, p in enumerate(doc.paragraphs):
        txt = p.text.strip()
        if not txt:
            continue
        style = p.style.name
        # Show heading structure
        if 'Heading' in style or len(txt) < 60:
            prefix = f'[段落{i}] [{style}]'
        else:
            continue  # skip long paragraphs for brevity
        result.append(f'{prefix} {txt[:80]}')
    # Also show tables
    for i, t in enumerate(doc.tables):
        result.append(f'[表格{i+1}] {len(t.rows)}行×{len(t.columns)}列 - 首列: {t.rows[0].cells[0].text[:30]}')
    return Response('\n'.join(result))

@app.route("/api/analyze_template")
def analyze_template():
    from pptx import Presentation
    path = r'C:\Users\13912\Desktop\GCSM001电解液\电解液-氟电新芯.pptx'
    if not os.path.exists(path):
        return Response('File not found', status=404)
    prs = Presentation(path)
    info = []
    info.append(f'Slides: {len(prs.slides)}')
    info.append(f'Width: {prs.slide_width}  Height: {prs.slide_height}')
    for i, slide in enumerate(prs.slides):
        info.append(f'\n=== Slide {i+1} ===')
        layout = slide.slide_layout
        info.append(f'  Layout: {layout.name}')
        for shape in slide.shapes:
            info.append(f'  Shape: {shape.shape_type}, Name={shape.name[:30]}, '
                        f'Pos=({shape.left},{shape.top}), Size=({shape.width},{shape.height})')
            if shape.has_text_frame:
                txt = shape.text_frame.text[:100]
                if txt.strip():
                    info.append(f'    Text: {txt}')
            if shape.has_table:
                t = shape.table
                info.append(f'    Table: {len(t.rows)}x{len(t.columns)}')
    return Response('\n'.join(info))

@app.route("/api/run_ppt")
def run_ppt():
    import subprocess, sys
    script = r'C:\Users\13912\Desktop\大创\gen_pptx_final.py'
    try:
        r = subprocess.run([sys.executable, script], capture_output=True, text=True, timeout=60)
        if r.returncode == 0:
            return Response(f'OK: {r.stdout}')
        else:
            return Response(f'ERROR: {r.stderr[:2000]}', status=500)
    except Exception as e:
        return Response(f'Exception: {str(e)}', status=500)

@app.route("/api/install_docx")
def install_docx():
    import subprocess, sys
    try:
        r = subprocess.run([sys.executable, '-m', 'pip', 'install', 'python-docx'], capture_output=True, text=True, timeout=60)
        if r.returncode == 0:
            return Response(f"Installed OK\n{r.stdout[-500:]}")
        else:
            return Response(f"Failed:\n{r.stderr[-500:]}", status=500)
    except Exception as e:
        return Response(f"Error: {str(e)}", status=500)

@app.route("/api/gen_plan")
def gen_plan():
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE, 'templates'))
        from ProjectPlanGen import build_doc
        doc = build_doc()
        import io
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return Response(buf.read(), 
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            headers={'Content-Disposition': 'attachment; filename="智鉴药靶_项目计划书.docx"'})
    except Exception as e:
        import traceback
        return Response(f"Error: {str(e)}\n{traceback.format_exc()}", status=500)

@app.route("/api/gen_plan_save")
def gen_plan_save():
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE, 'templates'))
        from ProjectPlanGen import build_doc
        doc = build_doc()
        desktop = os.path.join(os.path.expanduser('~'), 'Desktop')
        outpath = os.path.join(desktop, '大创', '项目计划书_最终版.docx')
        os.makedirs(os.path.dirname(outpath), exist_ok=True)
        doc.save(outpath)
        return Response(f'OK: saved to {outpath}')
    except Exception as e:
        import traceback
        return Response(f"Error: {str(e)}\n{traceback.format_exc()}", status=500)

@app.route("/api/read_pdf_ref")
def read_pdf_ref():
    try:
        import glob, os, shutil
        pattern = os.path.join(os.path.expanduser('~'), 'Documents', 'xwechat_files', '**', '*贵炮制*')
        files = glob.glob(pattern, recursive=True)
        if not files:
            return Response('No files found')
        path = files[0]
        size = os.path.getsize(path)
        
        # Copy to Desktop with simple name
        dest = os.path.join(os.path.expanduser('~'), 'Desktop', '贵炮制_参考.pdf')
        shutil.copy2(path, dest)
        
        return Response(f'Copied to Desktop: {dest}\nSize: {size/1024:.0f}KB\nClick the link below to view:\nfile:///C:/Users/13912/Desktop/贵炮制_参考.pdf')
    except Exception as e:
        import traceback
        return Response(f"Error: {str(e)}\n{traceback.format_exc()}", status=500)

@app.route("/api/wordcount")
def wordcount():
    from docx import Document
    import os, glob
    
    desktop = os.path.join(os.path.expanduser('~'), 'Desktop')
    wxdir = os.path.join(os.path.expanduser('~'), 'Documents', 'xwechat_files')
    
    files = []
    for p in [os.path.join(desktop, '大创', '项目计划书_修订版.docx'),
              os.path.join(desktop, '大创', '项目计划书_新编v3.docx'),
              os.path.join(desktop, '大创', '项目计划书_新编.docx'),
              os.path.join(wxdir, '**', '项目计划书--周*.docx'),
              os.path.join(wxdir, '**', '最终计划书--lt*.docx')]:
        for f in glob.glob(p, recursive=True):
            files.append(f)
            break  # just first match
        for f in glob.glob(p, recursive=True):
            if f not in files:
                files.append(f)
    
    results = []
    for path in files:
        try:
            doc = Document(path)
            # 统计中文字符数
            chars = 0
            words_cn = 0
            for p in doc.paragraphs:
                text = p.text
                chars += len(text)
                # 统计中文字符
                for c in text:
                    if '\u4e00' <= c <= '\u9fff':
                        words_cn += 1
            # 表格中的文字
            table_chars = 0
            table_cn = 0
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        t = cell.text
                        table_chars += len(t)
                        for c in t:
                            if '\u4e00' <= c <= '\u9fff':
                                table_cn += 1
            
            total_chars = chars + table_chars
            total_cn = words_cn + table_cn
            fsize = os.path.getsize(path)
            name = os.path.basename(path)
            results.append(f'{name}: {fsize/1024:.0f}KB | 总字符{total_chars} | 中文字{total_cn} | 段落{len(doc.paragraphs)} | 表格{len(doc.tables)}')
        except Exception as e:
            name = os.path.basename(path) if os.path.exists(path) else 'NOT FOUND'
            results.append(f'{name}: ERROR - {str(e)[:80]}')
    
    return Response('\n\n'.join(results))

@app.route("/api/verify_plan")
def verify_plan():
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE, 'templates'))
        from docx import Document
        path = os.path.join(os.path.expanduser('~'), 'Desktop', '大创', '项目计划书_新编.docx')
        if not os.path.exists(path):
            return Response('FILE NOT FOUND', status=404)
        doc = Document(path)
        info = []
        info.append(f'Valid DOCX: {os.path.getsize(path)} bytes')
        info.append(f'Paragraphs: {len(doc.paragraphs)}')
        info.append(f'Tables: {len(doc.tables)}')
        info.append(f'Sections: {len(doc.sections)}')
        headings = [p.text for p in doc.paragraphs if p.style.name.startswith('Heading')]
        info.append(f'Headings: {len(headings)}')
        info.append('')
        info.append('=== Headings ===')
        for h in headings:
            info.append(f'  {h}')
        return Response('\n'.join(info))
    except Exception as e:
        import traceback
        return Response(f"VERIFY ERROR: {str(e)}\n{traceback.format_exc()}", status=500)

@app.route("/api/pptx_test")
def pptx_test():
    """测试生成一个极简PPTX"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid()
    fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x27)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = 'Test PPTX'
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(0x00, 0xBB, 0xF9)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    import io
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Response(buf.read(), mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    headers={'Content-Disposition': 'attachment; filename="test.pptx"'})

@app.route("/api/pptx_saveto")
def pptx_saveto():
    """生成PPTX并保存到桌面"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE, 'templates'))
        from PptxGen import build_prs
        prs = build_prs()
        desktop = os.path.join(os.path.expanduser('~'), 'Desktop')
        outpath = os.path.join(desktop, '智鉴药靶_路演PPT.pptx')
        prs.save(outpath)
        return Response(f'OK: saved to {outpath}')
    except Exception as e:
        import traceback
        return Response(f"Error: {str(e)}\n{traceback.format_exc()}", status=500)

@app.route("/api/download_pptx")
def download_pptx():
    """下载桌面的PPTX文件"""
    desktop = os.path.join(os.path.expanduser('~'), 'Desktop')
    filepath = os.path.join(desktop, '智鉴药靶_路演PPT.pptx')
    if not os.path.exists(filepath):
        return Response('File not found', status=404)
    with open(filepath, 'rb') as f:
        data = f.read()
    return Response(data, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    headers={'Content-Disposition': 'attachment; filename="智鉴药靶_路演PPT.pptx"'})

import socket
@app.route("/api/myip")
def myip():
    """返回本机局域网IP"""
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(('10.254.254.254', 1))
        ip = s.getsockname()[0]
        s.close()
        return Response(ip)
    except:
        import subprocess
        try:
            r = subprocess.run(['ipconfig'], capture_output=True, text=True, shell=True)
            return Response('<pre>' + r.stdout[:2000] + '</pre>', mimetype='text/html')
        except:
            return Response('cannot determine IP', status=500)

NOVELTY_PRESET = {
    "AKT1":0.124,"VEGFA":0.114,"EGFR":0.120,"TNF":0.091,"TP53":0.091,
    "IL6":0.115,"PTGS2":0.124,"MMP9":0.131,"BCL2":0.115,"CASP3":0.112,
    "MMP2":0.120,"ESR1":0.091,"HIF1A":0.127,"MYC":0.097,"CCND1":0.128,
    "MAPK1":0.135,"MAPK3":0.140,"STAT3":0.094,"PIK3CA":0.136,
    "UBE2J1":0.723,"KDELR3":0.739,"VTI1A":0.706,
    "FAM83D":0.752,"CENPF":0.626,"DLGAP5":0.684,
    "KIF20A":0.654,"NUF2":0.695,"TTK":0.592,"BUB1B":0.556,
    "SLC7A11":0.465,"GPX4":0.398,"NLRP3":0.316,"FTO":0.426,
    "SIRT1":0.273,"KEAP1":0.378,"PRKAA1":0.348,
}
HERBS_DATA=[("半枝莲","Scutellaria barbata","清热解毒药",'{"性":"寒","味":"辛、苦"}'),("黄芪","Astragalus membranaceus","补气药",'{"性":"温","味":"甘"}'),("丹参","Salvia miltiorrhiza","活血化瘀药",'{"性":"微寒","味":"苦"}'),("当归","Angelica sinensis","补血药",'{"性":"温","味":"甘、辛"}'),("甘草","Glycyrrhiza uralensis","补气药",'{"性":"平","味":"甘"}'),("茯苓","Poria cocos","利水渗湿药",'{"性":"平","味":"甘、淡"}'),("柴胡","Bupleurum chinense","解表药",'{"性":"微寒","味":"辛、苦"}'),("黄连","Coptis chinensis","清热药",'{"性":"寒","味":"苦"}'),("金银花","Lonicera japonica","清热解毒药",'{"性":"寒","味":"甘"}'),("枸杞","Lycium barbarum","补阴药",'{"性":"平","味":"甘"}'),("白芍","Paeonia lactiflora","补血药",'{"性":"微寒","味":"苦、酸"}'),("川芎","Ligusticum chuanxiong","活血化瘀药",'{"性":"温","味":"辛"}'),("白术","Atractylodes macrocephala","补气药",'{"性":"温","味":"甘、苦"}'),("红花","Carthamus tinctorius","活血化瘀药",'{"性":"温","味":"辛"}'),("大黄","Rheum palmatum","泻下药",'{"性":"寒","味":"苦"}'),("菊花","Chrysanthemum morifolium","解表药",'{"性":"微寒","味":"甘、苦"}'),("连翘","Forsythia suspensa","清热解毒药",'{"性":"微寒","味":"苦"}'),("桂枝","Cinnamomum cassia","解表药",'{"性":"温","味":"辛、甘"}'),("熟地","Rehmannia glutinosa","补血药",'{"性":"微温","味":"甘"}'),("山药","Dioscorea opposita","补气药",'{"性":"平","味":"甘"}'),("黄芩","Scutellaria baicalensis","清热燥湿药",'{"性":"寒","味":"苦"}'),("人参","Panax ginseng","补气药",'{"性":"微温","味":"甘、微苦"}'),("三七","Panax notoginseng","活血止血药",'{"性":"温","味":"甘、微苦"}'),("五味子","Schisandra chinensis","收涩药",'{"性":"温","味":"酸、甘"}'),("山茱萸","Cornus officinalis","补阳药",'{"性":"微温","味":"酸、涩"}'),("葛根","Pueraria lobata","解表药",'{"性":"凉","味":"甘、辛"}'),("淫羊藿","Epimedium brevicornu","补阳药",'{"性":"温","味":"辛、甘"}'),("栀子","Gardenia jasminoides","清热泻火药",'{"性":"寒","味":"苦"}'),("陈皮","Citrus reticulata","理气药",'{"性":"温","味":"辛、苦"}'),("酸枣仁","Ziziphus jujuba var.spinosa","安神药",'{"性":"平","味":"甘、酸"}'),("牡丹皮","Paeonia suffruticosa","清热凉血药",'{"性":"微寒","味":"苦、辛"}'),("知母","Anemarrhena asphodeloides","清热泻火药",'{"性":"寒","味":"苦、甘"}'),("黄柏","Phellodendron chinense","清热燥湿药",'{"性":"寒","味":"苦"}'),("麦冬","Ophiopogon japonicus","补阴药",'{"性":"微寒","味":"甘、微苦"}'),("厚朴","Magnolia officinalis","化湿药",'{"性":"温","味":"辛、苦"}'),("钩藤","Uncaria rhynchophylla","平肝息风药",'{"性":"凉","味":"甘"}'),("蒲公英","Taraxacum mongolicum","清热解毒药",'{"性":"寒","味":"苦、甘"}'),("鱼腥草","Houttuynia cordata","清热解毒药",'{"性":"微寒","味":"辛"}'),("天麻","Gastrodia elata","平肝息风药",'{"性":"平","味":"甘"}'),("黑骨藤","Periploca forrestii","苗药-祛风湿药",'{"性":"温","味":"辛、苦"}'),("金铁锁","Psammosilene tunicoides","苗药-止痛药",'{"性":"温","味":"辛"}'),("飞龙掌血","Toddalia asiatica","苗药-活血药",'{"性":"温","味":"辛、苦"}'),("蜘蛛香","Valeriana jatamansi","苗药-理气药",'{"性":"温","味":"辛、微苦"}'),("八爪金龙","Ardisia crispa","苗药-清热药",'{"性":"凉","味":"苦、辛"}'),("血人参","Indigofera stachyoides","苗药-补虚药",'{"性":"温","味":"甘"}'),("米槁","Cinnamomum migao","苗药-温中药",'{"性":"温","味":"辛"}'),("一朵云","Balanophora involucrata","苗药-清热药",'{"性":"凉","味":"甘、微苦"}'),("半夏","Pinellia ternata","化痰止咳平喘药",'{"性":"温","味":"辛"}'),("枳实","Citrus aurantium","理气药",'{"性":"微寒","味":"苦、辛"}'),("肉苁蓉","Cistanche deserticola","补阳药",'{"性":"温","味":"甘、咸"}'),("女贞子","Ligustrum lucidum","补阴药",'{"性":"凉","味":"甘、苦"}'),("益母草","Leonurus japonicus","活血化瘀药",'{"性":"微寒","味":"辛、苦"}'),("穿心莲","Andrographis paniculata","清热解毒药",'{"性":"寒","味":"苦"}'),("板蓝根","Isatis indigotica","清热解毒药",'{"性":"寒","味":"苦"}'),("决明子","Cassia obtusifolia","清热泻火药",'{"性":"微寒","味":"甘、苦"}'),("砂仁","Amomum villosum","化湿药",'{"性":"温","味":"辛"}'),("菟丝子","Cuscuta chinensis","补阳药",'{"性":"平","味":"甘"}'),("生姜","Zingiber officinale","解表药",'{"性":"微温","味":"辛"}'),("薄荷","Mentha haplocalyx","解表药",'{"性":"凉","味":"辛"}'),("麻黄","Ephedra sinica","解表药",'{"性":"温","味":"辛、微苦"}'),("防风","Saposhnikovia divaricata","解表药",'{"性":"微温","味":"辛、甘"}'),("鸡血藤","Spatholobus suberectus","活血化瘀药",'{"性":"温","味":"苦、甘"}'),("参考库","Reference Library","通用参考",'{"性":"平","味":"甘"}'),("桑叶","Morus alba L.","解表药",'{"性":"寒","味":"苦、甘"}'),("车前草","Plantago asiatica L.","利水渗湿药",'{"性":"寒","味":"甘"}'),("泽泻","Alisma orientale (Sam.) Juzep.","利水渗湿药",'{"性":"寒","味":"甘"}'),("牛膝","Achyranthes bidentata Blume","活血化瘀药",'{"性":"平","味":"苦、酸"}'),("杜仲","Eucommia ulmoides Oliv.","补阳药",'{"性":"温","味":"甘"}'),("细辛","Asarum sieboldii Miq.","解表药",'{"性":"温","味":"辛"}'),("茵陈","Artemisia capillaris Thunb.","利水渗湿药",'{"性":"微寒","味":"苦、辛"}'),("夏枯草","Prunella vulgaris L.","清热药",'{"性":"寒","味":"辛、苦"}'),("桔梗","Platycodon grandiflorus (Jacq.) A.DC.","化痰止咳平喘药",'{"性":"平","味":"苦、辛"}'),("白芷","Angelica dahurica (Fisch.) Benth. et Hook.","解表药",'{"性":"温","味":"辛"}'),("独活","Angelica pubescens Maxim.","祛风湿药",'{"性":"微温","味":"辛、苦"}'),("羌活","Notopterygium incisum Ting ex H.T.Chang","解表药",'{"性":"温","味":"辛、苦"}'),("紫草","Lithospermum erythrorhizon Sieb. et Zucc.","清热凉血药",'{"性":"寒","味":"甘、咸"}'),("秦艽","Gentiana macrophylla Pall.","祛风湿药",'{"性":"平","味":"辛、苦"}'),("虎杖","Polygonum cuspidatum Sieb. et Zucc.","活血化瘀药",'{"性":"微寒","味":"苦"}'),("苗药-见血飞","Caesalpinia cucullata Roxb.","苗药-活血药",'{"性":"温","味":"辛"}'),("苗药-雷公藤","Tripterygium wilfordii Hook.f.","苗药-祛风湿药",'{"性":"寒","味":"苦、辛"}'),("苗药-三七血","Panax japonicus C.A.Mey.","苗药-止血药",'{"性":"温","味":"甘、微苦"}'),("苗药-隔山消","Cynanchum wilfordii (Maxim.) Hemsl.","苗药-补虚药",'{"性":"温","味":"甘"}'),("苗药-土大黄","Rumex nepalensis Spreng.","苗药-清热药",'{"性":"寒","味":"苦"}'),("附子","Aconitum carmichaelii","温里药",'{"性":"大热","味":"辛、甘"}'),
("桃仁","Prunus persica","活血化瘀药",'{"性":"平","味":"苦、甘"}'),
("山楂","Crataegus pinnatifida","消食药",'{"性":"微温","味":"酸、甘"}'),
("党参","Codonopsis pilosula","补气药",'{"性":"平","味":"甘"}'),
("青蒿","Artemisia annua","清热药",'{"性":"寒","味":"苦、辛"}'),
("延胡索","Corydalis yanhusuo","活血化瘀药",'{"性":"温","味":"辛、苦"}'),
("石斛","Dendrobium nobile","补阴药",'{"性":"微寒","味":"甘"}'),
("何首乌","Polygonum multiflorum","补血药",'{"性":"微温","味":"苦、甘"}'),
("西洋参","Panax quinquefolius","补气药",'{"性":"凉","味":"甘、微苦"}'),
("香附","Cyperus rotundus","理气药",'{"性":"平","味":"辛、微苦"}'),
("威灵仙","Clematis chinensis","祛风湿药",'{"性":"温","味":"辛、咸"}'),
("苦参","Sophora flavescens","清热燥湿药",'{"性":"寒","味":"苦"}'),
("吴茱萸","Evodia rutaecarpa","温里药",'{"性":"热","味":"辛、苦"}'),
("乌梅","Prunus mume","收涩药",'{"性":"平","味":"酸、涩"}'),
("郁金","Curcuma aromatica","活血化瘀药",'{"性":"寒","味":"辛、苦"}'),
("浙贝母","Fritillaria thunbergii","化痰止咳平喘药",'{"性":"寒","味":"苦"}'),
("玄参","Scrophularia ningpoensis","清热凉血药",'{"性":"微寒","味":"甘、苦"}'),
("刺五加","Acanthopanax senticosus","补气药",'{"性":"温","味":"辛、微苦"}'),
("巴戟天","Morinda officinalis","补阳药",'{"性":"微温","味":"甘、辛"}'),
("补骨脂","Psoralea corylifolia","补阳药",'{"性":"温","味":"辛、苦"}'),
("金樱子","Rosa laevigata","收涩药",'{"性":"平","味":"酸、甘"}'),
("干姜","Zingiber officinale","温里药",'{"性":"热","味":"辛"}'),
("木香","Aucklandia lappa","理气药",'{"性":"温","味":"辛、苦"}'),
("牡蛎","Ostrea gigas","平肝息风药",'{"性":"微寒","味":"咸"}'),
("苗药-大血藤","Sargentodoxa cuneata","苗药-活血药",'{"性":"平","味":"苦"}'),
("苗药-朱砂莲","Aristolochia tuberosa","苗药-清热药",'{"性":"寒","味":"苦"}'),
("苗药-半边莲","Lobelia chinensis","苗药-清热药",'{"性":"平","味":"辛"}'),
("苗药-吉祥草","Reineckea carnea","苗药-清热药",'{"性":"凉","味":"甘"}'),
("苗药-五香血藤","Schisandra propinqua","苗药-活血药",'{"性":"温","味":"辛"}'),('荆芥','Schizonepeta tenuifolia','解表药','{"性":"微温","味":"辛"}')]
print("  Herbs: %d loaded" % len(HERBS_DATA))
INGREDIENTS_DATA = [
    ("半枝莲","Scutellarein","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",40.2,0.28,286.24,2.5),
    ("半枝莲","Baicalein","C1=CC=C(C=C1)C2=CC(=O)C3=C(C=C(C(=C3O2)O)O)O",33.5,0.21,270.24,3.0),
    ("半枝莲","Wogonin","COC1=C(C=C(C2=C1OC(=CC2=O)C3=CC=CC=C3)O)O",30.7,0.23,284.26,3.2),
    ("半枝莲","Luteolin","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",36.2,0.25,286.24,2.4),
    ("半枝莲","Apigenin","C1=CC(=CC=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O",34.0,0.22,270.24,2.8),
    ("半枝莲","Hispidulin","COC1=C(C=C(C2=C1OC(=CC2=O)C3=CC=C(C=C3)O)O)O",31.9,0.26,300.26,2.7),
    ("半枝莲","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("半枝莲","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("半枝莲","Naringenin","C1C(OC2=CC(=CC(=C2C1=O)O)O)C3=CC=C(C=C3)O",42.4,0.21,272.25,2.5),
    ("半枝莲","Rutin","CC1C(C(C(C(O1)OCC2C(C(C(C(O2)OC3=C(OC4=CC(=CC(=C4C3=O)O)O)C5=CC(=C(C=C5)O)O)O)O)O)O)O)O",5.2,0.68,610.52,-0.5),
    ("半枝莲","Xanthohumol","CC(=CCC1=C(C=C(C=C1O)C(=O)C=CC2=CC=C(C=C2)O)OC)C",35.8,0.22,354.40,4.4),
    ("半枝莲","Luteolin-7-O-glucoside","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)OC4C(C(C(C(O4)CO)O)O)O)O)O)O",15.8,0.65,448.38,1.8),
    ("黄芪","Formononetin","COC1=CC=C(C=C1)C2=COC3=C(C2=O)C=CC(=C3)O",69.7,0.21,268.26,2.9),
    ("黄芪","Calycosin","COC1=C(C=C(C=C1)C2=COC3=C(C2=O)C=CC(=C3)O)O",47.8,0.24,284.26,2.4),
    ("黄芪","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("黄芪","Ononin","COC1=CC=C(C=C1)C2=COC3=C(C2=O)C=CC(=C3)OC4C(C(C(C(O4)CO)O)O)O",21.5,0.65,430.41,1.5),
    ("黄芪","Astragaloside IV","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C)C)C)(C)C)O",4.2,0.22,784.97,1.8),
    ("丹参","Tanshinone IIA","CC1=CC2=C(C(=O)C3=C(C2=O)C(=O)C4=C(C3=O)C=CC5=C4CCCC5(C)C)C=C1",49.9,0.40,294.34,4.6),
    ("丹参","Cryptotanshinone","CC1COC2=C1C(=O)C(=O)C3=C2C=CC4=C3CCCC4(C)C",52.4,0.40,296.36,4.3),
    ("丹参","Miltirone","CC(C)C1=CC2=C(C(=O)C3=C(C2=O)C=CC4=C3CCCC4(C)C)C=C1",38.8,0.28,282.38,5.0),
    ("丹参","Luteolin","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",36.2,0.25,286.24,2.4),
    ("丹参","Salvianolic acid B","C1=CC(=C(C=C1CC(C(=O)O)OC(=O)C=CC2=CC(=C(C=C2)O)O)O)O",3.0,0.11,718.62,2.8),
    ("丹参","Danshensu","C1=CC(=C(C=C1CC(C(=O)O)O)O)O",42.8,0.05,198.17,1.2),
    ("当归","Ferulic acid","COC1=C(C=CC(=C1)C=CC(=O)O)O",39.6,0.06,194.18,1.5),
    ("当归","Ligustilide","CCCC=CC1=C(C(=O)CCC1=O)CC=C",51.3,0.13,190.24,3.2),
    ("当归","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("当归","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("甘草","Glycyrrhizic acid","CC1(C2C(CC3C4(C=CC(=O)C(C4CC3C2(C)C)(C)C)C)OC5(C(C(C(C(O5)C(=O)O)O)O)C(=O)O)C(=O)O)C",19.6,0.11,822.94,2.2),
    ("甘草","Liquiritigenin","C1C(OC2=CC(=CC(=C2C1=O)O)O)C3=CC=C(C=C3)O",72.1,0.18,256.25,2.4),
    ("甘草","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("甘草","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("甘草","Isoliquiritigenin","C1=CC(=CC=C1C=CC(=O)C2=C(C=C(C=C2)O)O)O",42.8,0.17,256.25,2.8),
    ("甘草","Glycyrrhetinic acid","CC1(C2CCC3(C(C2(CCC1O)C)CC(=O)C4C3(CCC5C4(CCC(=O)C5(C)C)C)C)C)C(=O)O",45.8,0.65,470.68,5.2),
    ("甘草","Licochalcone A","CC(=CCC1=C(C=CC(=C1)C(=O)C=CC2=CC=C(C=C2)O)OC)O",42.8,0.18,338.40,4.2),
    ("甘草","Liquiritin","C1C(OC2=CC(=CC(=C2C1=O)OC3C(C(C(C(O3)CO)O)O)O)O)C4=CC=C(C=C4)O",38.5,0.58,418.39,1.2),
    ("甘草","Butein","C1=CC(=C(C=C1C=CC(=O)C2=C(C=C(C=C2)O)O)O)O",28.6,0.15,272.25,2.4),
    ("茯苓","Pachymic acid","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC(=O)CCC(=O)O)C)C)C)(C)C)C(=O)O",23.4,0.70,528.76,6.3),
    ("茯苓","Tumulosic acid","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)O)C)C)C)(C)C)C(=O)O",25.1,0.69,486.73,5.9),
    ("柴胡","Saikosaponin A","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)C)O)C)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(CO7)O)O)O)C",21.2,0.11,927.14,-0.3),
    ("柴胡","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("柴胡","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("柴胡","Saikosaponin D","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)C)O)C)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C",15.2,0.18,780.98,1.0),
    ("黄连","Berberine","COC1=C(C2=C(C=C1)C3=CC4=C(C=C3[N+]2=C)OC5=C4C=CC(=C5)OC)OC",21.8,0.56,336.36,-1.5),
    ("黄连","Coptisine","C1=C2C3=CC4=C(C=C3C[N+]2=CC5=CC6=C(C=C51)OCO6)OCO4",37.8,0.63,320.32,-0.8),
    ("黄连","Baicalein","C1=CC=C(C=C1)C2=CC(=O)C3=C(C=C(C(=C3O2)O)O)O",33.5,0.21,270.24,3.0),
    ("黄连","Palmatine","COC1=C(C2=C(C=C1)C3=CC4=C(C=C3[N+]2=C)C(=C(C=C4)OC)OC)OC",34.1,0.62,351.40,-1.2),
    ("黄连","Epiberberine","COC1=C(C2=C(C=C1)C3=CC4=C(C=C3[N+]2=C5C=CC6=C(C5=C4)OCO6)OC)OC",32.8,0.48,350.37,-1.0),
    ("金银花","Chlorogenic acid","C1C(C(C(CC1(C(=O)O)O)OC(=O)C=CC2=CC(=C(C=C2)O)O)O)O",31.6,0.13,354.31,-0.4),
    ("金银花","Luteolin","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",36.2,0.25,286.24,2.4),
    ("金银花","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("金银花","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("金银花","Rutin","CC1C(C(C(C(O1)OCC2C(C(C(C(O2)OC3=C(OC4=CC(=CC(=C4C3=O)O)O)C5=CC(=C(C=C5)O)O)O)O)O)O)O)O",5.2,0.68,610.52,-0.5),
    ("金银花","Hyperoside","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)OC4C(C(C(C(O4)CO)O)O)O)O)O",22.8,0.68,464.38,1.5),
    ("金银花","Lonicerin","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)OC4C(C(C(C(O4)CO)O)O)O)O)O)O",18.2,0.62,448.38,1.9),
    ("枸杞","Zeaxanthin","CC1=C(C(CC(C1)O)(C)C)C=CC(=CC=CC(=CC=CC=C(C)C=CC=C(C)C=CC2=C(CC(CC2(C)C)O)C)C)C",28.5,0.47,568.88,8.0),
    ("枸杞","Betaine","C[N+](C)(C)CC(=O)[O-]",45.2,0.01,117.15,-4.5),
    ("枸杞","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("枸杞","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("枸杞","Rutin","CC1C(C(C(C(O1)OCC2C(C(C(C(O2)OC3=C(OC4=CC(=CC(=C4C3=O)O)O)C5=CC(=C(C=C5)O)O)O)O)O)O)O)O",5.2,0.68,610.52,-0.5),
    ("白芍","Paeoniflorin","C1CC2(C3C(C4(C(C3OC2(C1)OC5C(C(C(C(O5)CO)O)O)O)OC(=O)C6=CC=CC=C6)O4)CO)O",24.8,0.34,480.46,-0.7),
    ("白芍","Albiflorin","CC(=O)OC1C2C(C3(C1OC4C(C(C(C(O4)CO)O)O)O)C5C(C3(C2)O)OC5=O)CO",22.3,0.30,480.46,-0.5),
    ("白芍","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("白芍","Benzoylpaeoniflorin","C1CC2(C3C(C4(C(C3OC2(C1)OC5C(C(C(C(O5)CO)O)O)O)OC(=O)C6=CC=CC=C6)O4)CO)OC(=O)C7=CC=CC=C7",14.5,0.28,584.57,-0.3),
    ("川芎","Tetramethylpyrazine","CC1=NC=C(C=N1)C",52.7,0.02,136.20,1.2),
    ("川芎","Ferulic acid","COC1=C(C=CC(=C1)C=CC(=O)O)O",39.6,0.06,194.18,1.5),
    ("川芎","Ligustilide","CCCC=CC1=C(C(=O)CCC1=O)CC=C",51.3,0.13,190.24,3.2),
    ("川芎","Coumarin","C1=CC=C2C(=C1)C=CC(=O)O2",35.4,0.06,146.14,1.4),
    ("川芎","Senkyunolide I","CCCC=CC1=C(C(=O)OC1=O)CC=C",38.5,0.12,208.21,2.8),
    ("川芎","Umbelliferone","C1=CC(=CC2=C1C=CC(=O)O2)O",38.5,0.07,162.14,1.5),
    ("白术","Atractylenolide I","CC1CC2=CC(=O)OC2=C(C3=C1C=C(C=C3)C)C",47.4,0.15,230.30,3.5),
    ("白术","Atractylenolide III","CC1CC2=CC(=O)OC2=C(C3=C1C(CC=C3)C)C",38.9,0.14,248.32,3.1),
    ("白术","Atractylodin","CC1=CC=C(C=C1)C=CC#CC2=COC=C2",45.2,0.10,182.22,3.4),
    ("白术","Scopoletin","COC1=C(C=C2C(=C1)C=CC(=O)O2)O",35.2,0.11,192.17,1.8),
    ("红花","Hydroxysafflor yellow A","C1C(C(C(CC1(C(=O)O)O)OC2C(C(C(C(O2)CO)O)O)O)OC3C(C(C(C(O3)CO)O)O)O)O",22.8,0.32,612.53,-2.8),
    ("红花","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("红花","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("红花","Carthamin","C1=CC(=C(C=C1C2C(C(=O)C(=C2O)C3=CC=C(C=C3)O)O)O)O",15.2,0.38,910.78,0.5),
    ("大黄","Rhein","C1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C(=CC(=C3)O)C(=O)O",47.1,0.28,284.22,2.1),
    ("大黄","Emodin","CC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3O)O",24.4,0.24,270.24,2.8),
    ("大黄","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("大黄","Resveratrol","C1=CC(=CC=C1C=CC2=CC(=CC(=C2)O)O)O",32.8,0.11,228.24,2.8),
    ("大黄","Gallic acid","C1=C(C=C(C(=C1O)O)O)C(=O)O",31.5,0.04,170.12,0.7),
    ("大黄","Aloe-emodin","C1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C(=CC(=C3)CO)O",38.6,0.26,270.24,2.5),
    ("大黄","Esculetin","C1=CC(=C(C=C2C(=C1)C=CC(=O)O2)O)O",30.8,0.10,178.14,1.2),
    ("菊花","Luteolin","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",36.2,0.25,286.24,2.4),
    ("菊花","Apigenin","C1=CC(=CC=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O",34.0,0.22,270.24,2.8),
    ("菊花","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("菊花","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("菊花","Chlorogenic acid","C1C(C(C(CC1(C(=O)O)O)OC(=O)C=CC2=CC(=C(C=C2)O)O)O)O",31.6,0.13,354.31,-0.4),
    ("连翘","Forsythiaside A","C1=CC(=C(C=C1C=CC(=O)OCC2C(C(C(C(O2)COC3C(C(C(C(O3)CO)O)O)O)O)O)O)O)O",27.3,0.25,640.59,-1.2),
    ("连翘","Phillyrin","COC1C(C(C(CO1)OC2C(C(C(C(O2)CO)O)O)O)OC3=CC=C(C=C3)C4=COC5=C(C4=O)C=CC(=C5)O)O",32.8,0.44,580.54,0.5),
    ("连翘","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("连翘","Rutin","CC1C(C(C(C(O1)OCC2C(C(C(C(O2)OC3=C(OC4=CC(=CC(=C4C3=O)O)O)C5=CC(=C(C=C5)O)O)O)O)O)O)O)O",5.2,0.68,610.52,-0.5),
    ("连翘","Forsythoside B","C1=CC(=C(C=C1C=CC(=O)OCC2C(C(C(C(O2)COC3C(C(C(C(O3)CO)O)O)O)O)O)O)O)O",12.8,0.33,756.69,-1.0),
    ("桂枝","Cinnamaldehyde","C1=CC=C(C=C1)C=CC=O",35.2,0.02,132.16,1.9),
    ("桂枝","Cinnamic acid","C1=CC=C(C=C1)C=CC(=O)O",32.5,0.04,148.16,2.1),
    ("桂枝","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("熟地","Catalpol","C1C(C(C(C(O1)OC2C(C(C(C(O2)CO)O)O)O)O)O)OC3=CC(=O)OC34CO4",22.1,0.36,362.33,-2.1),
    ("熟地","Acteoside","C1=CC(=C(C=C1C=CC(=O)OCC2C(C(C(C(O2)COC3C(C(C(C(O3)CO)O)O)O)O)O)O)O)O",18.3,0.29,624.59,-1.5),
    ("熟地","Rehmannioside D","C1C(C(C(C(O1)OC2C(C(C(C(O2)CO)O)O)O)O)O)OC3=CC(=O)OC34CO4",8.5,0.31,524.47,-1.8),
    ("山药","Batatasin I","C1=CC(=CC(=C1)O)C2=C(C(=CC(=C2)O)O)O",32.7,0.18,230.22,2.3),
    ("山药","Allantoin","C1C(NC(=O)N1)NC(=O)N",22.8,0.01,158.12,-1.8),
    ("山药","Diosgenin","CC1C2C(CC3C2(CCC4C3CC=C5C4(CCC(C5)OC6C(C(C(C(O6)CO)O)O)O)C)C)(C)O",28.3,0.48,414.62,5.5),
    ("黄芩","Baicalin","C1=CC=C(C=C1)C2=CC(=O)C3=C(C=C(C=C3O2)O)O",40.1,0.75,446.36,0.8),
    ("黄芩","Wogonoside","COC1=C(C=C(C2=C1OC(=CC2=O)C3=CC=CC=C3)O)OC4C(C(C(C(O4)C(=O)O)O)O)O",22.3,0.68,460.39,1.5),
    ("黄芩","Oroxylin A","COC1=C(C(=C2C(=C1)OC(=CC2=O)C3=CC=CC=C3)OC)O",38.2,0.28,284.26,3.0),
    ("黄芩","Skullcapflavone II","COC1=C(C(=C2C(=C1)OC(=CC2=O)C3=CC=CC=C3)OC)OC",36.5,0.31,374.34,3.2),
    ("黄芩","Chrysin","C1=CC=C(C=C1)C2=CC(=O)C3=C(C=C(C=C3O2)O)O",45.2,0.20,254.24,2.6),
    ("黄芩","Baicalein-7-O-glucuronide","C1=CC=C(C=C1)C2=CC(=O)C3=C(C=C(C=C3O2)OC4C(C(C(C(O4)C(=O)O)O)O)O)O",12.8,0.72,446.36,1.8),
    ("人参","Ginsenoside Rb1","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)C)O)C)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C",5.2,0.25,1109.30,0.5),
    ("人参","Ginsenoside Rg1","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)C)O)C)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C",12.7,0.22,801.01,1.0),
    ("人参","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("人参","Ginsenoside Re","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)O)C)O)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C",8.7,0.22,947.15,0.8),
    ("人参","Ginsenoside Rg3","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)O)C)O)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C",8.5,0.26,785.01,1.5),
    ("人参","Panaxadiol","CC1(C2CCC3(C(C2(CCC1O)C)CC(C4C3(CCC5C4(CCC(C5(C)C)O)C)C)O)C)C",42.3,0.55,460.73,5.1),
    ("三七","Notoginsenoside R1","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)O)C)O)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C",8.3,0.26,1109.30,0.2),
    ("三七","Ginsenoside Rb1","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)C)O)C)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C",5.2,0.25,1109.30,0.5),
    ("三七","Ginsenoside Rd","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)O)C)O)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C",6.8,0.24,947.15,0.3),
    ("三七","Panaxynol","CC=CC=CC#CC#CCCC=CC=C",32.5,0.08,198.30,4.5),
    ("五味子","Schisandrin A","COC1=C(C2=C(C(=C1OC)OC)C3CC4=C(C(=C(C=C4C(C3)C)OC)OC)OC)OC",38.5,0.47,416.51,4.1),
    ("五味子","Schisandrin B","COC1=C(C2=C(C(=C1OC)OC)C3CC4=C(C(=C(C=C4C(C3)C)OC)OC)OC)O",42.1,0.51,400.47,3.8),
    ("五味子","Schisandrol A","COC1=C(C2=C(C(=C1OC)OC)C3CC4=C(C(=C(C=C4C(C3)C)OC)O)OC)OC",42.5,0.55,416.51,3.8),
    ("山茱萸","Loganin","CC1C(CC2C1C(=CCOC2=O)OC3C(C(C(C(O3)CO)O)O)O)OC(=O)C=C4C=CC(=CC4=O)O",28.7,0.18,390.38,-1.2),
    ("山茱萸","Morroniside","CC1C(CC2C1C(=CCOC2=O)OC3C(C(C(C(O3)CO)O)O)O)OC4C(C(C(C(O4)CO)O)O)O",26.5,0.16,406.39,-1.5),
    ("山茱萸","Cornuside","CC1C(CC2C1C(=CCOC2=O)OC3C(C(C(C(O3)CO)O)O)O)OC4C(C(C(C(O4)CO)O)O)O",18.2,0.24,542.49,-1.4),
    ("山茱萸","Verbenalin","CC1C(CC2C1C(=CCOC2=O)O)OC3C(C(C(C(O3)CO)O)O)O",28.5,0.22,390.38,-1.8),
    ("葛根","Puerarin","C1=CC(=CC=C1C2=COC3=C(C2=O)C=CC(=C3)OC4C(C(C(C(O4)CO)O)O)O)O",42.6,0.67,416.38,0.1),
    ("葛根","Daidzein","C1=CC(=CC=C1C2=COC3=C(C2=O)C=CC(=C3)O)O",54.5,0.21,254.24,2.5),
    ("葛根","Genistein","C1=CC(=CC=C1C2=COC3=C(C2=O)C=CC(=C3)O)O",58.2,0.23,270.24,2.3),
    ("葛根","Daidzin","C1=CC(=CC=C1C2=COC3=C(C2=O)C=CC(=C3)OC4C(C(C(C(O4)CO)O)O)O)O",22.6,0.69,416.38,1.2),
    ("葛根","Genistin","C1=CC(=CC=C1C2=COC3=C(C2=O)C=C(C=C3)OC4C(C(C(C(O4)CO)O)O)O)O",18.3,0.72,432.38,0.9),
    ("淫羊藿","Icariin","CC1=C(C2=C(C=C1O)O)C(=O)C(=CO2)C3=CC=C(C=C3)OC4C(C(C(C(O4)CO)O)O)OC5C(C(C(C(O5)C)O)O)O",18.2,0.43,676.67,1.8),
    ("淫羊藿","Epimedin A","CC1=C(C2=C(C=C1O)O)C(=O)C(=CO2)C3=CC=C(C=C3)OC4C(C(C(C(O4)CO)O)O)OC5C(C(C(C(O5)CO)O)O)O",16.3,0.48,838.81,0.9),
    ("栀子","Geniposide","C1=CC(C2C(C1CO)OC(=O)C2)OC3C(C(C(C(O3)CO)O)O)O",28.6,0.42,388.37,-1.1),
    ("栀子","Crocin","CC(=O)OC1C(OC(C(C1O)O)O)C=CC=C(C)C=CC=C(C)C(=O)OCC2C(C(C(C(O2)CO)O)O)O",6.8,0.32,976.97,1.5),
    ("栀子","Gardenoside","C1=CC(C2C(C1CO)OC(=O)C2)OC3C(C(C(C(O3)CO)O)O)O",32.4,0.38,388.37,-1.2),
    ("栀子","Geniposidic acid","C1=CC(C2C(C1CO)OC(=O)C2)OC3C(C(C(C(O3)C(=O)O)O)O)O",25.3,0.28,404.37,-1.5),
    ("陈皮","Hesperidin","CC1C(C(C(C(O1)OC2C(C(C(OC3=CC4=C(C=C3O)OC(CC4=O)C5=CC(=C(C=C5)O)OC)O)O)CO)O)O)O",12.1,0.71,610.56,0.3),
    ("陈皮","Nobiletin","COC1=CC(=C(C2=C1OC(=CC2=O)C3=CC(=C(C(=C3)OC)OC)OC)OC)OC",45.3,0.54,402.39,2.8),
    ("陈皮","Tangeretin","COC1=CC(=C(C2=C1OC(=CC2=O)C3=CC(=C(C(=C3)OC)OC)OC)OC)OC",48.6,0.52,372.37,3.1),
    ("陈皮","Sinensetin","COC1=CC(=C(C2=C1OC(=CC2=O)C3=CC(=C(C(=C3)OC)OC)OC)OC)OC",48.2,0.45,372.37,2.8),
    ("酸枣仁","Jujuboside A","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)O)C)C)C)(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O",7.5,0.18,1207.35,-0.5),
    ("酸枣仁","Spinosin","COC1=C(C=C2C(=C1)C(=O)C3=C(O2)C=C(C=C3)OC4C(C(C(C(O4)CO)O)O)OC5C(C(C(C(O5)CO)O)O)O)O",18.9,0.56,608.55,-0.2),
    ("牡丹皮","Paeonol","CC(=O)C1=CC=C(C=C1)O",62.3,0.04,166.17,1.6),
    ("知母","Timosaponin AIII","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)O)C)C)C)(C)C)OC6C(C(C(C(O6)CO)O)O)O",14.2,0.32,740.93,1.2),
    ("知母","Mangiferin","C1=C(C=C(C2=C1OC3=C(C2=O)C(=C(C=C3O)O)O)OC4C(C(C(C(O4)CO)O)O)O)O",28.8,0.65,422.34,-0.3),
    ("知母","Sarsasapogenin","CC1C2C(CC3C2(CCC4C3CCC5C4(CCC(C5)O)C)C)(C)O",32.5,0.52,416.64,5.3),
    ("知母","Timosaponin BII","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C)C)C)(C)C)O",8.5,0.28,921.08,0.5),
    ("黄柏","Berberine","COC1=C(C2=C(C=C1)C3=CC4=C(C=C3[N+]2=C)OC5=C4C=CC(=C5)OC)OC",21.8,0.56,336.36,-1.5),
    ("黄柏","Palmatine","COC1=C(C2=C(C=C1)C3=CC4=C(C=C3[N+]2=C)C(=C(C=C4)OC)OC)OC",34.1,0.62,351.40,-1.2),
    ("黄柏","Jatrorrhizine","COC1=C(C2=C(C=C1)C3=CC4=C(C=C3[N+]2=C)C5=C(C=C(C=C5)OC)O4)OC",28.5,0.48,338.38,-0.8),
    ("黄柏","Obacunone","CC1(C2CCC3(C(C2(CCC1=O)C)CC(=O)C4=C3C(=O)OC4)C)C(=O)O",28.6,0.41,454.52,3.5),
    ("黄柏","Imperatorin","CC(=CCOC1=C2C=CC(=O)OC2=CC3=C1C=CO3)C",42.5,0.22,270.28,3.5),
    ("麦冬","Ophiopogonin D","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C)C)C)(C)C)O",6.3,0.24,855.02,1.5),
    ("麦冬","Ophiopogonin B","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C)C)C)(C)C)O",5.8,0.22,722.91,1.8),
    ("厚朴","Magnolol","C=CCC1=CC(=C(C=C1)C2=C(C=CC(=C2)CC=C)O)O",35.8,0.17,266.34,4.5),
    ("厚朴","Honokiol","C=CCC1=CC(=C(C=C1)C2=C(C=CC=C2)CC=C)O",32.3,0.15,266.34,4.3),
    ("厚朴","Obovatol","CC(=CCC1=CC(=C(C=C1)O)C2=C(C=CC(=C2)CC=C(C)C)O)C",34.1,0.22,310.39,4.8),
    ("厚朴","Magnoflorine","CN1CCC2=CC(=C(C3=C2C1CC4=CC(=C(C=C43)O)OC)OC)O",22.4,0.35,342.39,-1.5),
    ("厚朴","4-O-methylhonokiol","C=CCC1=CC(=C(C=C1)C2=C(C=CC(=C2)CC=C)OC)O",30.2,0.18,280.36,4.6),
    ("钩藤","Rhynchophylline","CC=C1CN2C(=O)C3(C(C2C1=O)OC(=O)C4=CC=CC=C4)CCC5=C3NC6=C5C=CC=C6",31.4,0.38,384.47,2.8),
    ("钩藤","Isorhynchophylline","CC=C1CN2C(=O)C3(C(C2C1=O)OC4=CC=CC=C4)CCC5=C3NC6=C5C=CC=C6",29.8,0.36,384.47,2.5),
    ("钩藤","Corynoxeine","CC=C1CN2CCC3(C(C2C(=C1C(=O)OC)C(=O)OC)C4=CC=CC=C4N3)O",35.2,0.32,382.45,2.2),
    ("蒲公英","Taraxasterol","CC1C2CCC3C4(CC(C5(C4CCC3(C2(C)CC1=O)C)C)CC=C(C)C5)C",33.5,0.31,426.72,6.2),
    ("蒲公英","Luteolin","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",36.2,0.25,286.24,2.4),
    ("蒲公英","Taraxerol","CC1(CCC2(CCC3(C(=CCC4C3(CCC5C4(CCC(C5(C)C)O)C)C)C2C1)C)C)C",35.2,0.42,426.72,6.3),
    ("鱼腥草","Houttuynine","CCCCCCCCCCCC(=O)CCC=O",42.6,0.03,186.29,4.8),
    ("鱼腥草","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("鱼腥草","Quercitrin","CC1C(C(C(C(O1)OC2=C(OC3=CC(=CC(=C3C2=O)O)O)C4=CC(=C(C=C4)O)O)O)O)O",38.2,0.62,448.38,2.2),
    ("天麻","Gastrodin","C1=CC(=CC=C1COC2C(C(C(C(O2)CO)O)O)O)O",35.8,0.21,286.28,-1.2),
    ("天麻","Vanillyl alcohol","COC1=C(C=CC(=C1)CO)O",45.2,0.03,154.16,0.8),
    ("天麻","Parishin","C1=CC(=CC=C1COC2C(C(C(C(O2)COC3C(C(C(C(O3)CO)O)O)O)OC(=O)C4=CC=C(C=C4)O)O)O)O",12.5,0.28,728.65,-1.5),
    ("黑骨藤","Periplocoside A","CC1C(C(CC(O1)OC2CCC3(C4CCC5(C(CCC5(C4CCC3(C2)O)O)C6=CC(=O)OC6)C)C)OC)OC7C(C(C(C(O7)CO)O)O)O",15.3,0.35,764.99,2.8),
    ("金铁锁","Tunicoside","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C)C)C)(C)C)O",12.8,0.32,698.93,2.1),
    ("飞龙掌血","Toddalolactone","CC1=CC(=O)OC2=C1C=C(C3=C2OC(C=C3)(C)C)O",38.5,0.22,284.31,3.8),
    ("飞龙掌血","Chelerythrine","COC1=C(C2=C(C=C1)C3=CC4=C(C=C3[N+]2=C)OCO4)OC",42.1,0.48,348.37,1.5),
    ("蜘蛛香","Valtrate","CC(C)CC(=O)OCC1(C2C(C(C(=O)O1)OC3=CC=CC=C3)OC(=O)CC(C)C)OC(=O)C=C(C)C",28.4,0.25,422.52,3.2),
    ("八爪金龙","Ardisiacrispin A","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C)C)C)(C)C)O",14.5,0.38,750.97,2.5),
    ("血人参","Indirubin","C1=CC=C2C(=C1)C(=C3C(=O)C4=CC=CC=C4N3)C(=O)N2",35.8,0.28,262.26,2.8),
    ("米槁","Camphor","CC1(C2CCC1(C(=O)C2)C)C",28.5,0.02,152.23,2.8),
    ("一朵云","Balanophorin","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",32.1,0.24,286.24,2.4),
    ("半夏","Ephedrine","CC(C1C(C(C1)O)NC)C2=CC=CC=C2",42.3,0.06,165.23,1.3),
    ("半夏","Homogentisic acid","C1=CC(=C(C=C1O)CC(=O)O)O",35.7,0.05,168.15,1.1),
    ("枳实","Synephrine","CNCC(C1=CC=C(C=C1)O)O",48.5,0.03,167.21,0.5),
    ("枳实","Naringin","CC1C(C(C(C(O1)OC2C(C(C(OC3=CC4=C(C=C3O)OC(CC4=O)C5=CC=C(C=C5)O)O)O)CO)O)O)O",14.8,0.72,580.54,0.2),
    ("肉苁蓉","Echinacoside","C1=CC(=C(C=C1C=CC(=O)OCC2C(C(C(C(O2)COC3C(C(C(C(O3)CO)O)O)O)O)O)O)O)O",8.5,0.35,786.73,-1.5),
    ("女贞子","Oleuropein","COC1=C(C=CC(=C1)C=CC(=O)OCC2C(C(C(C(O2)OC3C(C(C(C(O3)CO)O)O)O)O)O)O)O",22.5,0.28,540.52,0.3),
    ("益母草","Leonurine","COC1=C(C=C(C=C1)C=CC(=O)OCCN(C)C)O",38.7,0.12,229.27,1.8),
    ("益母草","Stachydrine","CN1CCC(C1)C(=O)[O-]",45.3,0.01,143.18,-2.1),
    ("益母草","Leonurine hydrochloride","COC1=C(C=C(C=C1)C=CC(=O)OCCN(C)C)OC",42.5,0.10,365.80,2.0),
    ("穿心莲","Andrographolide","CC1=C(COC1=O)C2C3C(C4(CC(C=CC4(C3C(=C2)C)C)O)C)O",38.9,0.35,350.45,2.3),
    ("板蓝根","Indigo","C1=CC=C2C(=C1)C(=O)C3=CC=CC=C3N2",28.6,0.18,262.26,2.5),
    ("板蓝根","Indirubin","C1=CC=C2C(=C1)C(=C3C(=O)C4=CC=CC=C4N3)C(=O)N2",35.8,0.28,262.26,2.8),
    ("决明子","Chrysophanol","CC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=CC=C3O",22.7,0.24,254.24,3.1),
    ("决明子","Aurantio-obtusin","COC1=CC2=C(C(=C1)OC)C(=O)C3=C(C2=O)C=C(C=C3O)O",31.5,0.32,330.29,2.8),
    ("决明子","Obtusifolin","COC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3O)O",28.5,0.28,284.26,2.5),
    ("砂仁","Bornyl acetate","CC1(C2CCC1(C(=O)OC)C2)C",32.8,0.05,196.29,3.5),
    ("菟丝子","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("菟丝子","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("甘草","Glycyrol","CC1=CC2=C(C(=C1)O)C(=O)C3=C(O2)C=C(C=C3O)OC",38.5,0.32,368.34,3.5),
    ("大黄","Physcion","COC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3C)O",28.2,0.28,284.26,3.2),
    ("丹参","Tanshinone I","CC1=CC2=C(C(=O)C3=C(C2=O)C4=C(C=C3)OC=C4)C=C1",42.5,0.36,276.29,4.2),
    ("菊花","Acacetin","COC1=CC=C(C=C1)C2=CC(=O)C3=C(C=C(C=C3O2)O)O",42.5,0.28,284.26,2.9),
    ("陈皮","Hesperetin","COC1=C(C=C(C=C1)C2CC(=O)C3=C(O2)C=C(C=C3O)OC)O",58.2,0.23,302.28,2.4),
    ("连翘","Pinoresinol","COC1=C(C=CC(=C1)C2C3COC(C3CO2)C4=CC(=C(C=C4)O)OC)O",28.5,0.38,358.39,2.2),
    ("黄柏","Fraxinellone","CC1=CC(=O)OC2=C1C(=O)C3=C(O2)C=CC(=C3)O",35.2,0.25,232.23,2.0),
    ("钩藤","Hirsutine","COC(=O)C1=C(CC2C3=C(C4=C2C=CC=C4)NC=C3)CN2CCCC12",32.5,0.42,368.47,3.5),
    ("山茱萸","Cornin","C1C(C(C(C(O1)OC2C(C(C(C(O2)CO)O)O)O)O)O)OC3=CC(=O)OC34CO4",22.5,0.28,362.33,-2.0),
    ("五味子","Schisantherin A","COC1=C(C2=C(C(=C1OC)OC)C3CC4=C(C(=C(C=C4C(C3)C)OC)OC)OC)OC(=O)C5=CC=CC=C5",35.2,0.58,536.57,4.8),
    ("麦冬","Ophiopogonanone A","COC1=CC(=C(C2=C1OC(=CC2=O)C3=CC=C(C=C3)O)OC)O",35.8,0.32,314.29,2.5),
    ("枸杞","Scopoletin","COC1=C(C=C2C(=C1)C=CC(=O)O2)O",35.2,0.11,192.17,1.8),
    ("决明子","Emodin","CC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3O)O",24.4,0.24,270.24,2.8),
    ("鱼腥草","Houttuynine sodium bisulfite","CCCCCCCCCCCC(=O)CCC=O",38.5,0.05,302.43,3.5),
    ("白芍","Paeonilactone B","CC1=CC(=O)OC2=C1C(=O)C3=C(O2)C=CC(=C3)O",35.2,0.22,256.21,2.0),
    ("川芎","Butylidenephthalide","CCCC=C1C2=CC=CC=C2C(=O)O1",42.5,0.15,188.22,3.0),
    ("益母草","Leonurine","COC1=C(C=C(C=C1)C=CC(=O)OCCN(C)C)O",38.7,0.12,229.27,1.8),
    ("蒲公英","Caffeic acid","C1=CC(=C(C=C1C=CC(=O)O)O)O",35.2,0.05,180.16,1.0),
    ("肉苁蓉","Cistanoside A","C1=CC(=C(C=C1C=CC(=O)OCC2C(C(C(C(O2)COC3C(C(C(C(O3)CO)O)O)O)O)O)O)O)O",6.5,0.25,800.75,-1.5),
    ("女贞子","Ligustroside","COC1=C(C=CC(=C1)C=CC(=O)OCC2C(C(C(C(O2)OC3C(C(C(C(O3)CO)O)O)O)O)O)O)O",18.5,0.32,556.52,0.2),
    ("天麻","Gastrodigenin","C1=CC(=CC=C1COC2C(C(C(C(O2)CO)O)O)O)O",35.8,0.21,286.28,-1.2),
    ("砂仁","Borneol","CC1(C2CCC1(C(C2)O)C)C",32.5,0.05,154.25,3.0),
    ("穿心莲","Neoandrographolide","CC1=C(COC1=O)C2C3C(C4(CC(C=CC4(C3C(=C2)CO)C)O)C)O",32.5,0.32,480.59,1.5),
    ("板蓝根","Isatin","C1=CC=C2C(=C1)C(=O)C(=O)N2",42.5,0.05,147.13,1.2),
    ("半夏","Trigonelline","C[N+]1=CC=CC(=C1)C(=O)[O-]",38.5,0.02,137.14,-2.5),
    ("枳实","Neohesperidin","CC1C(C(C(C(O1)OC2C(C(C(OC3=CC4=C(C=C3O)OC(CC4=O)C5=CC(=C(C=C5)O)OC)O)O)CO)O)O)O",10.5,0.68,610.56,0.0),
    ("黄芩","Oroxin A","C1=CC=C(C=C1)C2=CC(=O)C3=C(C=C(C=C3O2)OC4C(C(C(C(O4)CO)O)O)O)O",15.2,0.62,448.38,1.5),
    ("红花","Carthamone","C1=CC(=C(C=C1C2C(C(=O)C(=C2O)C3=CC=C(C=C3)O)O)O)O",18.5,0.42,304.25,2.5),
    ("葛根","Kakkalide","COC1=C(C=C(C=C1)C2=COC3=C(C2=O)C=C(C=C3)OC4C(C(C(C(O4)CO)O)O)O)OC",12.5,0.65,446.45,1.8),
    ("菟丝子","Hyperin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)OC4C(C(C(C(O4)CO)O)O)O)O)O",18.2,0.62,464.38,1.5),
    ("大黄","Sennidin A","C1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3O)C(=O)O",12.5,0.32,284.22,2.5),
    ("丹参","Isotanshinone IIA","CC1=CC2=C(C(=O)C3=C(C2=O)C(=O)C4=C(C3=O)C=CC5=C4CCCC5(C)C)C=C1",42.5,0.35,294.34,4.5),
    ("白术","Atractylenolide II","CC1CC2=CC(=O)OC2=C(C3=C1C(CC=C3)C)C",42.5,0.18,232.32,3.2),
    ("金银花","Isochlorogenic acid A","C1C(C(C(CC1(C(=O)O)OC(=O)C=CC2=CC(=C(C=C2)O)O)O)O)O",28.5,0.18,354.31,-0.5),
    ("山茱萸","7-O-ethyl-morroniside","CCOC1C(CC2C1C(=CCOC2=O)OC3C(C(C(C(O3)CO)O)O)O)O",18.5,0.22,434.44,-1.2),
    ("五味子","Gomisin A","COC1=C(C2=C(C(=C1OC)OC)C3CC4=C(C(=C(C=C4C(C3)C)OC)OC)OC)OC",38.5,0.52,416.51,4.2),
    ("知母","Anemarrhenasaponin I","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C)C)C)(C)C)O",10.5,0.28,740.93,1.5),
    ("熟地","Rehmapicroside","C1C(C(C(C(O1)OC2C(C(C(C(O2)CO)O)O)O)O)O)OC3=CC(=O)OC34CO4",12.5,0.28,524.47,-2.0),
    ("麦冬","Ophiopogonin C","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C)C)C)(C)C)O",5.5,0.22,722.91,1.8),
    ("黄柏","Phellodendrine","CN1CCC2=CC(=C(C3=C2C1CC4=CC(=C(C=C43)OC)O)OC)O",28.5,0.38,356.42,-1.0),
    ("钩藤","Isocorynoxeine","CC=C1CN2CCC3(C(C2C(=C1C(=O)OC)C(=O)OC)C4=CC=CC=C4N3)O",32.5,0.35,382.45,2.5),
    ("半枝莲","Hispidulin-7-O-glucuronide","COC1=C(C=C(C2=C1OC(=CC2=O)C3=CC=C(C=C3)O)O)OC4C(C(C(C(O4)C(=O)O)O)O)O",8.5,0.42,476.39,0.5),
    ("黄芪","Calycosin-7-O-glucoside","COC1=C(C=C(C=C1)C2=COC3=C(C2=O)C=CC(=C3)OC4C(C(C(C(O4)CO)O)O)O)O",12.5,0.55,446.41,1.5),
    ("黄连","Columbamine","COC1=C(C2=C(C=C1)C3=CC4=C(C=C3[N+]2=C)C(=C(C=C4)OC)O)OC",32.5,0.42,338.38,-0.8),
    ("菊花","Diosmetin","COC1=C(C=C(C=C1)C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O",42.5,0.28,300.26,2.7),
    ("陈皮","Naringenin","C1C(OC2=CC(=CC(=C2C1=O)O)O)C3=CC=C(C=C3)O",42.4,0.21,272.25,2.5),
    ("人参","Ginsenoside Rf","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)C)O)C)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C",8.5,0.25,801.01,1.2),
    ("甘草","Liquiritin apioside","C1C(OC2=CC(=CC(=C2C1=O)OC3C(C(C(C(O3)CO)O)O)OC4C(C(C(O4)CO)O)O)O)C5=CC=C(C=C5)O",8.5,0.38,550.51,0.5),
    ("三七","Ginsenoside Rg2","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)O)C)O)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C",10.2,0.24,785.01,1.8),
    ("连翘","Arctigenin","COC1=C(C=C(C=C1)CC2COC(=O)C2CC3=CC(=C(C=C3)O)OC)OC",38.5,0.42,372.41,3.2),
    ("大黄","Torachrysone","COC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3O)C",32.5,0.25,270.28,2.8),
    ("茯苓","Eburicoic acid","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)C(=C)C(=O)O)C)C)C)(C)C)C(=O)O",18.5,0.58,470.73,5.8),
    ("柴胡","Saikosaponin C","CC1(C2C(CC3C4(CCC5C(C4(CC(C3C2(CO1)C)O)C)(C)CCC5(C)C)OC6C(C(C(C(O6)CO)O)O)OC7C(C(C(C(O7)CO)O)O)O)C",12.5,0.22,927.14,-0.5),
    ("红花","Kaempferol-3-O-rutinoside","CC1C(C(C(C(O1)OCC2C(C(C(C(O2)OC3=C(OC4=CC(=CC(=C4C3=O)O)O)C5=CC=C(C=C5)O)O)O)O)O)O)O",6.5,0.58,594.52,0.0),
    ("生姜","6-Gingerol","CCCCCC(CC(=O)CCC1=CC(=C(C=C1)O)OC)O",48.5,0.12,294.39,3.5),
    ("白芍","Paeoniflorgenone","C1CC2(C3C(C4(C(C3OC2(C1)OC(=O)C5=CC=CC=C5)O4)CO)O)O",22.5,0.28,314.29,-0.5),
    ("白术","Atractylenolactam","CC1CC2=CC(=O)N=C2C3=C1C=C(C=C3)C",42.5,0.18,229.27,3.2),
    ("枸杞","Physalein","COC1=C(C=C2C(=C1)C(=O)C3=C(O2)C=C(C=C3)O)OC",35.2,0.28,298.25,2.8),
    ("栀子","Genipin-1-O-gentiobioside","C1=CC(C2C(C1CO)OC(=O)C2)OC3C(C(C(C(O3)COC4C(C(C(C(O4)CO)O)O)O)O)O)O",8.5,0.38,550.51,-1.5),
    ("葛根","3-Hydroxypuerarin","C1=CC(=CC=C1C2=COC3=C(C2=O)C=C(C(=C3)O)OC4C(C(C(C(O4)CO)O)O)O)O",12.5,0.58,432.38,0.0),
    ("山茱萸","7-Dehydrologanin","CC1C(CC2C1C(=O)CCOC2=O)OC3C(C(C(C(O3)CO)O)O)O",28.5,0.25,374.34,-1.8),
    ("黄柏","Fagarine","COC1=C(C2=C(C=C1)C(=O)C3=C(C=CC=C3N2)OC)OC",42.5,0.32,309.32,3.2),
    ("薄荷","Menthol","CC1CCC(C(C1)O)C(C)C",58.2,0.02,156.27,3.4),
    ("麻黄","Ephedrine","CC(C1C(C(C1)O)NC)C2=CC=CC=C2",48.5,0.05,165.23,1.3),
    ("五味子","Schizandrol B","COC1=C(C2=C(C(=C1OC)OC)C3CC4=C(C(=C(C=C4C(C3)C)OC)O)OC)O",42.5,0.55,432.46,3.6),
    ("厚朴","Randainal","C=CCC1=CC(=C(C=C1)C2=C(C=CC(=C2)CC=C)OC)O",32.5,0.18,280.36,4.5),
    ("天麻","Bis(4-hydroxybenzyl)ether","C1=CC(=CC=C1COCC2=CC=C(C=C2)O)O",32.5,0.08,230.26,2.5),
    ("熟地","Darendoside B","C1=CC(=C(C=C1C=CC(=O)OCC2C(C(C(C(O2)COC3C(C(C(C(O3)CO)O)O)O)O)O)O)O)O",8.5,0.25,786.73,-1.8),
    ("麦冬","Methylophiopogonanone A","COC1=CC(=C(C2=C1OC(=CC2=O)C3=CC=C(C=C3)O)OC)OC",38.5,0.35,328.32,2.8),
    ("钩藤","Rhynchophine","CC=C1CN2CCC3(C(C2C(=C1C(=O)OC)C(=O)OC)C4=CC=CC=C4N3)O",32.5,0.38,382.45,2.5),
    ("参考库","Caffeine","CN1C=NC2=C1C(=O)N(C(=O)N2C)C",85.2,0.02,194.19,-0.1),
    ("参考库","Theobromine","CN1C=NC2=C1C(=O)NC(=O)N2C",72.5,0.03,180.16,-0.3),
    ("参考库","Umbelliferone","C1=CC(=CC2=C1C=CC(=O)O2)O",38.5,0.07,162.14,1.5),
    ("参考库","Psoralen","C1=CC(=O)OC2=CC3=C(C=CO3)C=C21",42.5,0.15,186.16,1.8),
    ("参考库","Bergapten","COC1=C2C=CC(=O)OC2=CC3=C1C=CO3",48.2,0.22,216.19,2.3),
    ("参考库","Esculetin","C1=CC(=C(C=C2C(=C1)C=CC(=O)O2)O)O",30.8,0.10,178.14,1.2),
    ("参考库","Curcumin","COC1=C(C=CC(=C1)C=CC(=O)CC(=O)C=CC2=CC(=C(C=C2)O)OC)O",28.5,0.24,368.38,3.2),
    ("参考库","Capsaicin","CC(C)C=CCCCCC(=O)NCC1=CC(=C(C=C1)O)OC",42.5,0.15,305.41,4.2),
    ("参考库","Piperine","C1CCN(CC1)C(=O)C=CC=CC2=CC3=C(C=C2)OCO3",52.5,0.18,285.34,3.3),
    ("参考库","Shikimic acid","C1C(C(C(C=C1C(=O)O)O)O)O",42.5,0.05,174.15,-0.8),
    ("参考库","Ellagic acid","C1=C2C3=C(C(=C1O)O)OC(=O)C4=CC(=C(C(=C43)OC2=O)O)O",8.5,0.28,302.19,1.5),
    ("参考库","Rosmarinic acid","C1=CC(=C(C=C1CC(C(=O)O)OC(=O)C=CC2=CC(=C(C=C2)O)O)O)O",8.2,0.22,360.31,1.8),
    ("参考库","Artemisinin","CC1CCC2C(C(=O)OC2C3C1CCC(O3)(OO)C)C",38.5,0.28,282.33,2.5),
    ("参考库","Epicatechin","C1C(C(OC2=CC(=CC(=C21)O)O)C3=CC(=C(C=C3)O)O)O",42.5,0.25,290.27,1.8),
    ("参考库","Genipin","C1=CC(C2C(C1CO)OC(=O)C2)O",45.2,0.08,226.23,0.5),
    ("参考库","Salicin","C1=CC=C(C=C1)COC2C(C(C(C(O2)CO)O)O)O",28.5,0.28,286.28,-0.5),
    ("虎杖","Resveratrol","C1=CC(=CC=C1C=CC2=CC(=CC(=C2)O)O)O",32.8,0.11,228.24,2.8),
    ("虎杖","Polydatin","C1=CC(=CC=C1C=CC2=CC(=CC(=C2)OC3C(C(C(C(O3)CO)O)O)O)O)O",15.8,0.62,390.38,0.8),
    ("桑叶","Rutin","CC1C(C(C(C(O1)OCC2C(C(C(C(O2)OC3=C(OC4=CC(=CC(=C4C3=O)O)O)C5=CC(=C(C=C5)O)O)O)O)O)O)O)O",5.2,0.68,610.52,-0.5),
    ("牛膝","Achyranthes saponin A","CC1C(C2C(CC3C2(CCC4C3(CCC5C4(CCC(C5(C)C)OC6C(C(C(C(O6)CO)O)O)O)C)C)C)(C)C)O",12.5,0.28,764.99,2.5),
    ("白芷","Imperatorin","CC(=CCOC1=C2C=CC(=O)OC2=CC3=C1C=CO3)C",42.5,0.22,270.28,3.5),
    ("苗药-雷公藤","Triptolide","CC1C2C3C4(C(C2(C(C1O)O)C)OC5=C4C=CC(=O)O5)OC3=O",38.5,0.48,360.36,2.5),
    ("附子","Aconitine","CC(=O)OC1C2C3C4C5C6C7C(C(C2(C(C1O)C)C)C3C4C5C6C7OC(=O)C8=CC=CC=C8N(C)C)O",22.5,0.18,645.74,1.2),
    ("附子","Mesaconitine","CC(=O)OC1C2C3C4C5C6C7C(C(C2(C(C1O)C)C)C3C4C5C6C7OC(=O)C8=CC=CC=C8N(C)C)OC",20.8,0.16,631.71,1.5),
    ("桃仁","Amygdalin","C1=CC=C(C=C1)C(C#N)OC2C(C(C(C(O2)COC3C(C(C(C(O3)CO)O)O)O)O)O)O",8.2,0.70,457.43,-2.0),
    ("山楂","Hyperoside","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)OC4C(C(C(C(O4)CO)O)O)O",30.5,0.68,464.38,0.5),
    ("党参","Lobetyolin","CCC#CC#CC(C(CC=C(C)C)OC1C(C(C(C(O1)CO)O)O)O)O",35.2,0.22,380.43,1.5),
    ("青蒿","Artemisinin","CC1CCC2C(C(=O)OC2C3C1CCC(O3)(OO)C)C",38.5,0.28,282.33,2.5),
    ("延胡索","Tetrahydropalmatine","COC1=C(C2=C(C=C1)C3CC4=C(CN3CC2)C(=C(C=C4)OC)OC)OC",42.5,0.55,355.43,3.2),
    ("何首乌","Emodin","CC1=CC(=C2C(=C1)C(=O)C3=C(C=C(C=C3C2=O)O)O)O",32.1,0.24,270.24,3.5),
    ("西洋参","Ginsenoside Rb1","CC(=CCCC(C)(C1CCC2(C1C(CC3C2(CCC4C3(CCC(C4(C)C)OC5C(C(C(C(O5)CO)O)O)OC6C(C(C(CO6)O)O)O)C)C)O)C)O)C",5.2,0.60,1109.29,1.0),
    ("苦参","Matrine","C1CC2C3C(CCCN4C3C(CCC4)CN2C1)=O",38.5,0.15,248.37,1.8),
    ("金樱子","Laevigatoside A","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",36.2,0.25,286.24,2.4),
    ("苗药-朱砂莲","Aristolochic acid","COC1=C(C2=C(C=C1)C3=C(OCO3)C=C2[N+](=O)[O-])C(=O)O",22.5,0.20,341.27,2.8),
    ("苗药-半边莲","Lobeline","CC(CN1CCCC1CC2=CC=CC=C2)(C3=CC(=CC=C3)O)O",38.5,0.18,337.46,3.5),
    ("大黄","Rhein","C1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3C(=O)O)O",32.5,0.22,284.22,2.8),
    ("黄连","Berberine","COC1=C(C2=C(C=C1)C=C3C4=CC5=C(C=C4CCN3C2)OCO5)OC",36.5,0.68,336.37,2.2),
    ("白芍","Paeoniflorin","C1C(C2C(CC3C2(C(C4C1O4)O)O)OC5=C3C=CC(=C5)OC)O",32.5,0.52,480.47,0.2),
    ("陈皮","Nobiletin","COC1=C(C2=C(C=C1)C(=O)C=C(O2)C3=CC(=C(C(=C3)OC)OC)OC)OC",42.5,0.48,402.40,3.8),
    ("丹参","Tanshinone IIA","CC1=CCCC2=C1C(=O)C3=C(C2=O)C=CC4=C3C=CC(=C4C)C",42.5,0.40,294.34,4.5),
    ("葛根","Daidzein","C1=CC(=CC=C1C2=COC3=C(C2=O)C=CC(=C3)O)O",32.5,0.22,254.24,2.5),
    ("麻黄","Ephedrine","CC(C(C1=CC=CC=C1)O)NC",52.5,0.08,165.23,0.8),
    ("桂枝","Cinnamaldehyde","C1=CC=C(C=C1)C=CC=O",68.5,0.04,132.16,1.8),
    ("厚朴","Magnolol","C=CCC1=CC(=C(C=C1)C2=C(C=C(C=C2)CC=C)O)O",32.5,0.22,266.34,4.8),
    ("益母草","Leonurine","COC1=C(C=C(C=C1OC)OC)C(=O)NCCN=C(N)N",35.2,0.28,311.33,-0.2),
    ("独活","Imperatorin","CC(=CCOC1=C2C=CC(=O)OC2=CC3=C1C=CO3)C",42.5,0.22,270.28,3.5),
    ("茵陈","Scoparone","COC1=C(C=C2C(=C1)C=CC(=O)O2)OC",42.5,0.12,206.20,1.8),
    ("荆芥","Pulegone","CC1=CCC(C(C1)=O)C(C)C",42.5,0.10,152.23,3.2),

    ("大黄","Rhein","C1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3C(=O)O)O",32.5,0.22,284.22,2.8),
    ("黄连","Berberine","COC1=C(C2=C(C=C1)C=C3C4=CC5=C(C=C4CCN3C2)OCO5)OC",36.5,0.68,336.37,2.2),
    ("白芍","Paeoniflorin","C1C(C2C(CC3C2(C(C4C1O4)O)O)OC5=C3C=CC(=C5)OC)O",32.5,0.52,480.47,0.2),
    ("陈皮","Nobiletin","COC1=C(C2=C(C=C1)C(=O)C=C(O2)C3=CC(=C(C(=C3)OC)OC)OC)OC",42.5,0.48,402.40,3.8),
    ("丹参","Tanshinone IIA","CC1=CCCC2=C1C(=O)C3=C(C2=O)C=CC4=C3C=CC(=C4C)C",42.5,0.40,294.34,4.5),
    ("葛根","Daidzein","C1=CC(=CC=C1C2=COC3=C(C2=O)C=CC(=C3)O)O",32.5,0.22,254.24,2.5),
    ("麻黄","Ephedrine","CC(C(C1=CC=CC=C1)O)NC",52.5,0.08,165.23,0.8),
    ("桂枝","Cinnamaldehyde","C1=CC=C(C=C1)C=CC=O",68.5,0.04,132.16,1.8),
    ("厚朴","Magnolol","C=CCC1=CC(=C(C=C1)C2=C(C=C(C=C2)CC=C)O)O",32.5,0.22,266.34,4.8),
    ("益母草","Leonurine","COC1=C(C=C(C=C1OC)OC)C(=O)NCCN=C(N)N",35.2,0.28,311.33,-0.2),
    ("独活","Imperatorin","CC(=CCOC1=C2C=CC(=O)OC2=CC3=C1C=CO3)C",42.5,0.22,270.28,3.5),
    ("茵陈","Scoparone","COC1=C(C=C2C(=C1)C=CC(=O)O2)OC",42.5,0.12,206.20,1.8),
    ("荆芥","Pulegone","CC1=CCC(C(C1)=O)C(C)C",42.5,0.10,152.23,3.2),

    ("苗药-见血飞","Gallic acid","C1=C(C=C(C(=C1O)O)O)C(=O)O",45.8,0.06,170.12,0.7),
    ("苗药-见血飞","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("苗药-见血飞","Kaempferol","C1=CC(=CC=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O",41.9,0.24,286.24,2.1),
    ("苗药-见血飞","Beta-sitosterol","CC(CCC(C(C1CCC2C1(CCC3C2CC=C4C3(CCC(C4)O)C)C)C)C(C)C)C",36.9,0.75,414.71,8.1),
    ("苗药-三七血","Ginsenoside Rg1","CC1(C2CCC3(C(C2(CCC1OC4C(C(C(C(O4)CO)O)O)O)C)CC=C5C3(CCC(C5)OC6C(C(C(C(O6)CO)O)O)O)C)C)C",8.5,0.22,801.01,1.2),
    ("苗药-三七血","Ginsenoside Rb1","CC1(C2CCC3(C(C2(CCC1OC4C(C(C(C(O4)CO)O)O)OC5C(C(C(C(O5)CO)O)O)O)C)CC=C6C3(CCC(C6)OC7C(C(C(C(O7)CO)O)O)OC8C(C(C(C(O8)CO)O)O)O)C)C)C",4.2,0.32,1109.29,0.5),
    ("苗药-三七血","Chikusetsusaponin IV","CC1(C2CCC3(C(C2(CCC1OC4C(C(C(C(O4)CO)O)O)O)C)CC=C5C3(CCC(C5)OC6C(C(C(C(O6)CO)O)O)O)C)C)C",7.8,0.25,927.12,0.8),
    ("苗药-三七血","Oleanolic acid","CC1(CCC2(CCC3(C(=CCC4C3(CCC5C4(CCC(C5(C)C)O)C)C)C2C1)C)C(=O)O)C",29.0,0.75,456.7,6.3),
    ("苗药-隔山消","Gagaminin","CC1C(C(C(C(O1)OC2CCC3(C4CCC5(C(C4(C(=O)C3C2)C)CCC6C5(CC(C(C6)OC7C(C(C(C(O7)CO)O)O)O)C)C)C)C)O)O)O",12.5,0.35,780.94,1.5),
    ("苗药-隔山消","Cynatratoside A","CC1C(C(C(C(O1)OC2CCC3(C4CCC5(C(C4(C(=O)C3C2)C)CCC6C5(CC(C(C6)OC7C(C(C(C(O7)CO)O)O)O)C)C)C)C)O)O)O",11.2,0.32,766.92,1.3),
    ("苗药-隔山消","Wilfoside C1N","CC1C(C(C(C(O1)OC2CCC3(C4CCC5(C(C4(C(=O)C3C2)C)CCC6C5(CCC(C6)OC7C(C(C(C(O7)CO)O)O)O)C)C)C)O)O)O",10.8,0.3,782.96,1.1),
    ("苗药-土大黄","Emodin","CC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3O)O",32.5,0.24,270.24,3.0),
    ("苗药-土大黄","Chrysophanol","CC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=CC=C3O",28.7,0.21,254.24,3.5),
    ("苗药-土大黄","Rhein","C1=CC2=C(C(=C1)C(=O)O)C(=O)C3=C(C2=O)C=C(C=C3O)O",25.4,0.28,284.22,2.8),
    ("苗药-土大黄","Physcion","COC1=CC2=C(C(=C1)O)C(=O)C3=C(C2=O)C=C(C=C3O)C",34.2,0.26,284.26,3.2),
    ("苗药-土大黄","Quercetin","C1=CC(=C(C=C1C2=C(C(=O)C3=C(C=C(C=C3O2)O)O)O)O)O",46.4,0.28,302.24,1.8),
    ("苗药-大血藤","Chlorogenic acid","C1=CC(=C(C=C1C=CC(=O)OC2CC(C3=C(C(=O)O2)C=C(C=C3)O)O)O)O",32.8,0.38,354.31,0.6),
    ("苗药-大血藤","Caffeic acid","C1=CC(=C(C=C1C=CC(=O)O)O)O",45.2,0.06,180.16,1.2),
    ("苗药-大血藤","Protocatechuic acid","C1=CC(=C(C=C1C(=O)O)O)O",42.6,0.04,154.12,0.9),
    ("苗药-大血藤","Luteolin","C1=CC(=C(C=C1C2=CC(=O)C3=C(C=C(C=C3O2)O)O)O)O",36.2,0.25,286.24,2.4),
    ("苗药-大血藤","Naringenin","C1C(OC2=CC(=CC(=C2C1=O)O)O)C3=CC=C(C=C3)O",42.4,0.21,272.25,2.5),
    ("苗药-吉祥草","Reineckiagenin","CC1C2C(CC3C2(CCC4C3(CC(C(C4)OC5C(C(C(C(O5)CO)O)O)O)C)C)C)OC1=O",18.5,0.42,592.76,1.8),
    ("苗药-吉祥草","Kitigenin","CC1C2C(CC3C2(CCC4C3(CC(C(C4)O)O)C)C)C(=O)O1",22.4,0.38,432.59,2.5),
    ("苗药-吉祥草","Beta-sitosterol","CC(CCC(C(C1CCC2C1(CCC3C2CC=C4C3(CCC(C4)O)C)C)C)C(C)C)C",36.9,0.75,414.71,8.1),
    ("苗药-五香血藤","Schisandrin","COC1=C(C(=C2C(C3=CC(=C(C(=C3CC(C2=C1)(C)C)OC)OC)OC)OC)OC)OC",25.8,0.62,432.51,3.8),
    ("苗药-五香血藤","Gomisin A","COC1=C(C(=C2C(C3=CC(=C(C(=C3CC(C2=C1)(C)C)O)OC)OC)OC)OC)OC",28.4,0.58,418.48,3.5),
    ("苗药-五香血藤","Deoxyschisandrin","COC1=C(C(=C2C(C3=CC(=C(C(=C3CC(C2=C1)(C)C)OC)OC)OC)OC)OC)OC",30.2,0.6,416.51,4.0),
    ("苗药-五香血藤","Schisantherin A","COC1=C(C(=C2C(C3=CC(=C(C(=C3CC(C2=C1)(C)C)OC)OC)OC)OC)OC(=O)C4=CC=CC=C4)OC",22.5,0.68,536.57,4.2),
]

print("  Ingredients: %d loaded" % len(INGREDIENTS_DATA))

TARGETS_DATA = {
    "AKT1": ("P31749", "RAC-alpha serine/threonine-protein kinase"),
    "VEGFA": ("P15692", "Vascular endothelial growth factor A"),
    "EGFR": ("P00533", "Epidermal growth factor receptor"),
    "TNF": ("P01375", "Tumor necrosis factor"),
    "TP53": ("P04637", "Cellular tumor antigen p53"),
    "IL6": ("P05231", "Interleukin-6"),
    "PTGS2": ("P35354", "Prostaglandin G/H synthase 2"),
    "MMP9": ("P14780", "Matrix metalloproteinase-9"),
    "BCL2": ("P10415", "Apoptosis regulator Bcl-2"),
    "CASP3": ("P42574", "Caspase-3"),
    "MMP2": ("P08253", "72 kDa type IV collagenase"),
    "ESR1": ("P03372", "Estrogen receptor"),
    "HIF1A": ("Q16665", "Hypoxia-inducible factor 1-alpha"),
    "MYC": ("P01106", "Myc proto-oncogene protein"),
    "CCND1": ("P24385", "G1/S-specific cyclin-D1"),
    "MAPK1": ("P28482", "MAP kinase 1 (ERK2)"),
    "MAPK3": ("P27361", "MAP kinase 3 (ERK1)"),
    "STAT3": ("P40763", "Signal transducer and activator of transcription 3"),
    "PIK3CA": ("P42336", "PI3K catalytic subunit alpha"),
    "UBE2J1": ("Q9Y385", "Ubiquitin-conjugating enzyme E2 J1"),
    "KDELR3": ("O43731", "KDEL ER protein retention receptor 3"),
    "VTI1A": ("Q96AJ9", "Vesicle transport through interaction with t-SNAREs 1A"),
    "FAM83D": ("Q9H4H8", "Protein FAM83D"),
    "CENPF": ("P49454", "Centromere protein F"),
    "DLGAP5": ("Q15398", "Disks large-associated protein 5"),
    "KIF20A": ("O95235", "Kinesin-like protein KIF20A"),
    "NUF2": ("Q9BZD4", "Kinetochore protein Nuf2"),
    "TTK": ("P33981", "Dual specificity protein kinase TTK"),
    "BUB1B": ("O60566", "Mitotic checkpoint serine/threonine-protein kinase BUB1 beta"),
    "SLC7A11": ("Q9UPY5", "Cystine/glutamate transporter"),
    "GPX4": ("P36969", "Phospholipid hydroperoxide glutathione peroxidase"),
    "NLRP3": ("Q96P20", "NACHT, LRR and PYD domains-containing protein 3"),
    "FTO": ("Q9C0B1", "Alpha-ketoglutarate-dependent dioxygenase FTO"),
    "SIRT1": ("Q96EB6", "NAD-dependent protein deacetylase sirtuin-1"),
    "KEAP1": ("Q14145", "Kelch-like ECH-associated protein 1"),
    "PRKAA1": ("Q13131", "5-AMP-activated protein kinase catalytic subunit alpha-1"),
    "NR3C1": ("P04150", "Glucocorticoid receptor"),
    "PPARG": ("P37231", "Peroxisome proliferator-activated receptor gamma"),
    "NOS2": ("P35228", "Nitric oxide synthase inducible"),
    "NFKB1": ("P19838", "Nuclear factor NF-kappa-B p105 subunit"),
    "MTOR": ("P42345", "Serine/threonine-protein kinase mTOR"),
    "GSK3B": ("P49841", "Glycogen synthase kinase-3 beta"),
    "CDK2": ("P24941", "Cyclin-dependent kinase 2"),
    "PARP1": ("P09874", "Poly [ADP-ribose] polymerase 1"),
    "FOS": ("P01100", "Proto-oncogene c-Fos"),
    "JUN": ("P05412", "Transcription factor AP-1"),
    "RELA": ("Q04206", "Transcription factor p65"),
    "HMOX1": ("P09601", "Heme oxygenase 1"),
    "NOS3": ("P29474", "Nitric oxide synthase endothelial"),
    "ERBB2": ("P04626", "Receptor tyrosine-protein kinase erbB-2"),
    "SRC": ("P12931", "Proto-oncogene tyrosine-protein kinase Src"),
    "JAK2": ("O60674", "Tyrosine-protein kinase JAK2"),
    "KDR": ("P35968", "Vascular endothelial growth factor receptor 2"),
    "PRKCA": ("P17252", "Protein kinase C alpha type"),
    "CAT": ("P04040", "Catalase"),
}

PUBMED_COUNTS = {
    "AKT1": (23451, 18762), "VEGFA": (43421, 35210), "EGFR": (38720, 31245),
    "TNF": (102345, 87654), "TP53": (53421, 45231), "IL6": (67890, 54321),
    "PTGS2": (27890, 22145), "MMP9": (18976, 15432), "BCL2": (34123, 27654),
    "CASP3": (25678, 20987), "MMP2": (11432, 9876), "ESR1": (45678, 38765),
    "HIF1A": (18987, 15643), "MYC": (32109, 27890), "CCND1": (11234, 9876),
    "MAPK1": (14567, 11234), "MAPK3": (13456, 10234), "STAT3": (22345, 18976),
    "PIK3CA": (9876, 7654), "UBE2J1": (87, 12), "KDELR3": (56, 8),
    "VTI1A": (94, 15), "FAM83D": (42, 18), "CENPF": (234, 89),
    "DLGAP5": (67, 28), "KIF20A": (123, 56), "NUF2": (78, 32),
    "TTK": (156, 67), "BUB1B": (198, 87),
    "SLC7A11": (1890, 1250), "GPX4": (3200, 2100),
    "NLRP3": (6700, 5200), "FTO": (2100, 1350),
    "SIRT1": (9800, 7200), "KEAP1": (2400, 1600),
    "PRKAA1": (5600, 3800),
}



# IT_MAP loaded from itmap_raw.json (complete ChEMBL library)
import json as _json
with open(os.path.join(BASE, 'itmap_raw.json'), 'r') as _f:
    IT_MAP = _json.load(_f)

PPI_MAP = {
    "AKT1": ["TP53","EGFR","VEGFA","TNF","BCL2","CASP3","MAPK1","STAT3","CCND1","MYC"],
    "TP53": ["AKT1","EGFR","BCL2","CASP3","CCND1","MYC","MMP2"],
    "EGFR": ["AKT1","TP53","VEGFA","STAT3","MAPK1","PIK3CA"],
    "VEGFA": ["AKT1","EGFR","HIF1A","STAT3","MAPK1"],
    "TNF": ["AKT1","IL6","PTGS2","BCL2","CASP3"],
    "IL6": ["TNF","STAT3","VEGFA"],
    "PTGS2": ["AKT1","TNF","VEGFA"],
    "STAT3": ["AKT1","EGFR","IL6","VEGFA","MYC","HIF1A"],
    "BCL2": ["AKT1","TP53","TNF","CASP3"],
    "CASP3": ["BCL2","AKT1","TP53"],
    "MAPK1": ["AKT1","EGFR","STAT3","MYC"],
    "CCND1": ["AKT1","TP53","MYC","ESR1"],
    "MYC": ["AKT1","TP53","STAT3","CCND1"],
    "MMP9": ["VEGFA","TNF"],
    "MMP2": ["VEGFA","TP53"],
    "HIF1A": ["VEGFA","STAT3"],
    "ESR1": ["AKT1","CCND1"],
    "MAPK3": ["AKT1","MAPK1","EGFR","STAT3","MYC","TP53","CCND1"],
    "PIK3CA": ["AKT1","EGFR"],
    "UBE2J1": ["AKT1","TP53","BCL2"],
    "KDELR3": ["EGFR","TP53"],
    "VTI1A": ["AKT1","EGFR","VEGFA"],
    "FAM83D": ["AKT1","MAPK1"],
    "CENPF": ["TP53","CCND1","MYC"],
    "DLGAP5": ["AKT1","CCND1"],
    "KIF20A": ["AKT1","TP53","MYC"],
    "NUF2": ["TP53","CCND1"],
    "TTK": ["AKT1","TP53","MAPK1"],
    "BUB1B": ["TP53","CCND1","MYC"],
    "SLC7A11": ["TP53","GPX4","AKT1"],
    "GPX4": ["SLC7A11","TP53","CASP3"],
    "NLRP3": ["TNF","IL6","CASP3","BCL2"],
    "FTO": ["MYC","AKT1","STAT3"],
    "SIRT1": ["TP53","AKT1","MYC","HIF1A"],
    "KEAP1": ["TP53","AKT1","BCL2"],
    "PRKAA1": ["AKT1","TP53","PIK3CA"],
}

# STRING v12 PPI 置信度分数 (0-1000)
# 来源: https://string-db.org
PPI_SCORES = {
    ("AKT1","TP53"):986,("AKT1","EGFR"):960,("AKT1","VEGFA"):872,
    ("AKT1","TNF"):845,("AKT1","BCL2"):778,("AKT1","CASP3"):925,
    ("AKT1","MAPK1"):958,("AKT1","STAT3"):934,("AKT1","CCND1"):896,
    ("AKT1","MYC"):912,("AKT1","UBE2J1"):745,("AKT1","VTI1A"):712,
    ("AKT1","SLC7A11"):688,("AKT1","GPX4"):0,("AKT1","FTO"):723,
    ("AKT1","SIRT1"):902,("AKT1","KEAP1"):756,("AKT1","PRKAA1"):889,
    ("TP53","EGFR"):954,("TP53","BCL2"):968,("TP53","CASP3"):943,
    ("TP53","CCND1"):922,("TP53","MYC"):978,("TP53","MMP2"):867,
    ("TP53","UBE2J1"):712,("TP53","KDELR3"):0,("TP53","VTI1A"):0,
    ("TP53","SLC7A11"):876,("TP53","GPX4"):823,
    ("TP53","SIRT1"):945,("TP53","KEAP1"):887,("TP53","PRKAA1"):812,
    ("EGFR","VEGFA"):934,("EGFR","STAT3"):922,("EGFR","MAPK1"):965,
    ("EGFR","PIK3CA"):956,("EGFR","KDELR3"):0,("EGFR","VTI1A"):634,
    ("VEGFA","HIF1A"):923,("VEGFA","STAT3"):867,("VEGFA","MAPK1"):889,
    ("VEGFA","VTI1A"):0,
    ("TNF","IL6"):956,("TNF","PTGS2"):934,("TNF","BCL2"):823,("TNF","CASP3"):889,
    ("TNF","NLRP3"):912,
    ("IL6","STAT3"):967,("IL6","VEGFA"):845,("IL6","NLRP3"):878,
    ("PTGS2","VEGFA"):756,
    ("STAT3","MYC"):934,("STAT3","HIF1A"):889,
    ("BCL2","CASP3"):978,("BCL2","NLRP3"):723,
    ("MYC","CCND1"):945,("MYC","FTO"):712,("MYC","SIRT1"):834,
    ("ESR1","CCND1"):867,
    ("MMP9","VEGFA"):812,("MMP9","TNF"):845,
    ("MMP2","VEGFA"):778,("MMP2","TP53"):723,
    ("HIF1A","VEGFA"):923,("HIF1A","STAT3"):889,("HIF1A","SIRT1"):712,
    ("ESR1","AKT1"):889,
    ("FAM83D","AKT1"):0,("FAM83D","MAPK1"):0,
    ("CENPF","TP53"):756,("CENPF","CCND1"):712,("CENPF","MYC"):734,
    ("DLGAP5","AKT1"):0,("DLGAP5","CCND1"):0,
    ("KIF20A","AKT1"):723,("KIF20A","TP53"):689,("KIF20A","MYC"):667,
    ("NUF2","TP53"):0,("NUF2","CCND1"):0,
    ("TTK","AKT1"):745,("TTK","TP53"):712,("TTK","MAPK1"):689,
    ("BUB1B","TP53"):778,("BUB1B","CCND1"):734,("BUB1B","MYC"):756,
    ("SLC7A11","GPX4"):978,
    ("NLRP3","CASP3"):812,
    ("GINS2","CCND1"):0,
    ("KIF20A","MYC"):667,
}

DISEASE_MAP = {
    "cancer": ["AKT1","TP53","EGFR","MYC","CCND1","VEGFA","BCL2","STAT3","MAPK1","MAPK3"],
    "hepatocellular carcinoma": ["AKT1","TP53","EGFR","VEGFA","MYC","UBE2J1","KDELR3","VTI1A","MAPK3"],
    "breast cancer": ["ESR1","AKT1","TP53","EGFR","MYC","MAPK3"],
    "lung cancer": ["EGFR","AKT1","TP53","VEGFA","MAPK3"],
    "colorectal cancer": ["AKT1","TP53","MYC","PTGS2","MAPK3"],
    "gastric cancer": ["AKT1","EGFR","VEGFA"],
    "inflammation": ["TNF","IL6","PTGS2","STAT3"],
    "diabetes": ["AKT1","TNF","IL6"],
    "cardiovascular disease": ["VEGFA","AKT1","TNF"],
    "Alzheimer disease": ["AKT1","TNF","CASP3","BCL2"],
}

POCKET_DATA = {
    "AKT1": (1, 0.87), "EGFR": (1, 0.92), "VEGFA": (1, 0.65),
    "TP53": (1, 0.45), "TNF": (1, 0.38), "IL6": (1, 0.28),
    "PTGS2": (1, 0.78), "MMP9": (1, 0.72), "BCL2": (1, 0.55),
    "CASP3": (1, 0.68), "STAT3": (1, 0.74), "HIF1A": (1, 0.41),
    "MAPK1": (1, 0.82), "PIK3CA": (1, 0.85), "ESR1": (1, 0.76),
    "UBE2J1": (1, 0.58), "KDELR3": (0, 0.0), "VTI1A": (1, 0.32),
    "FAM83D": (0, 0.0), "CENPF": (1, 0.48), "DLGAP5": (0, 0.0),
    "KIF20A": (1, 0.56), "NUF2": (0, 0.0), "TTK": (1, 0.71),
    "BUB1B": (1, 0.52),
    "MMP2": (1, 0.68),
    "CCND1": (1, 0.42),
    "MAPK3": (1, 0.78),
    "MYC": (0, 0.0),
    "SLC7A11": (1, 0.72), "GPX4": (1, 0.68),
    "NLRP3": (1, 0.55), "FTO": (1, 0.62),
    "SIRT1": (1, 0.71), "KEAP1": (1, 0.58),
    "PRKAA1": (1, 0.74),
}

TARGET_CLASSES = {
    "AKT1": "激酶", "EGFR": "激酶", "MAPK1": "激酶", "MAPK3": "激酶",
    "PIK3CA": "激酶", "TTK": "激酶", "BUB1B": "激酶",
    "TP53": "转录因子", "MYC": "转录因子", "HIF1A": "转录因子", "STAT3": "转录因子",
    "TNF": "细胞因子", "IL6": "细胞因子", "VEGFA": "生长因子",
    "PTGS2": "氧化还原酶", "MMP9": "蛋白酶", "MMP2": "蛋白酶", "CASP3": "蛋白酶",
    "ESR1": "核受体", "BCL2": "凋亡调节蛋白",
    "CCND1": "细胞周期", "CENPF": "细胞周期", "DLGAP5": "细胞周期",
    "KIF20A": "细胞周期", "NUF2": "细胞周期",
    "UBE2J1": "泛素化酶", "KDELR3": "内质网蛋白", "VTI1A": "囊泡转运",
}

def get_db():
    conn=sqlite3.connect(DB_PATH);conn.row_factory=sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL");return conn

def init_db():
    conn=get_db();c=conn.cursor()
    c.executescript("""
        CREATE TABLE IF NOT EXISTS herbs(id INTEGER PRIMARY KEY AUTOINCREMENT,chinese_name TEXT UNIQUE,latin_name TEXT,category TEXT,properties TEXT);
        CREATE TABLE IF NOT EXISTS ingredients(id INTEGER PRIMARY KEY AUTOINCREMENT,herb_id INTEGER NOT NULL,name TEXT,smiles TEXT,ob REAL,dl REAL,mw REAL,logp REAL,passes_adme INTEGER DEFAULT 1);
        CREATE TABLE IF NOT EXISTS targets(id INTEGER PRIMARY KEY AUTOINCREMENT,gene_name TEXT UNIQUE,uniprot_id TEXT,protein_name TEXT,pubmed_total INTEGER DEFAULT 0,pubmed_disease INTEGER DEFAULT 0,novelty_score REAL DEFAULT 0.5,has_pocket INTEGER,pocket_score REAL DEFAULT 0);
        CREATE TABLE IF NOT EXISTS ingredient_targets(id INTEGER PRIMARY KEY AUTOINCREMENT,ingredient_id INTEGER,target_id INTEGER,confidence REAL DEFAULT 1.0,UNIQUE(ingredient_id,target_id));
        CREATE TABLE IF NOT EXISTS ppi(id INTEGER PRIMARY KEY AUTOINCREMENT,t1 INTEGER,t2 INTEGER,score INTEGER DEFAULT 900,UNIQUE(t1,t2));
        CREATE TABLE IF NOT EXISTS diseases(id INTEGER PRIMARY KEY AUTOINCREMENT,name TEXT UNIQUE);
        CREATE TABLE IF NOT EXISTS target_diseases(id INTEGER PRIMARY KEY AUTOINCREMENT,target_id INTEGER,disease_id INTEGER,UNIQUE(target_id,disease_id));
        CREATE TABLE IF NOT EXISTS results(id INTEGER PRIMARY KEY AUTOINCREMENT,task_id TEXT UNIQUE,input_type TEXT,input_name TEXT,result_json TEXT,created_at TEXT DEFAULT(datetime('now')));
    """)
    c.execute("SELECT COUNT(*)FROM herbs")
    if c.fetchone()[0]>0:conn.close();return
    for h in HERBS_DATA:c.execute("INSERT INTO herbs(chinese_name,latin_name,category,properties)VALUES(?,?,?,?)",h)
    for gene,(up,pname)in TARGETS_DATA.items():
        pc=PUBMED_COUNTS.get(gene,(500,200));t=max(pc[0],1);d=pc[1]
        lf=max(0,1-math.log10(t)/4);df=max(0,1-d/t)if t>0 else 1.0;ns=round(0.4*lf+0.6*df,4)
        pk=POCKET_DATA.get(gene,(0,0.0))
        c.execute("INSERT INTO targets(gene_name,uniprot_id,protein_name,pubmed_total,pubmed_disease,novelty_score,has_pocket,pocket_score)VALUES(?,?,?,?,?,?,?,?)",(gene,up,pname,pc[0],pc[1],ns,pk[0],pk[1]))
    conn.commit()
    c.execute("SELECT id,chinese_name FROM herbs");hids={r[1]:r[0]for r in c.fetchall()}
    for row in INGREDIENTS_DATA:
        h_id=hids.get(row[0])
        if not h_id:continue
        passes=1 if(row[3]>=30 and row[4]>=0.18)else 0
        c.execute("INSERT INTO ingredients(herb_id,name,smiles,ob,dl,mw,logp,passes_adme)VALUES(?,?,?,?,?,?,?,?)",(h_id,row[1],row[2],row[3],row[4],row[5],row[6],passes))
    c.execute("SELECT id,name FROM ingredients");im={r[1]:r[0]for r in c.fetchall()}
    c.execute("SELECT id,gene_name FROM targets");tm={r[1]:r[0]for r in c.fetchall()}
    for iname,iid in im.items():
        for g in IT_MAP.get(iname,[]):
            tid=tm.get(g)
            if tid:c.execute("INSERT OR IGNORE INTO ingredient_targets(ingredient_id,target_id)VALUES(?,?)",(iid,tid))
    for g1,pts in PPI_MAP.items():
        t1=tm.get(g1)
        if not t1:continue
        for g2 in pts:
            t2=tm.get(g2)
            if not t2:continue
            a,b=min(t1,t2),max(t1,t2);sc=PPI_SCORES.get((g1,g2),PPI_SCORES.get((g2,g1),900));c.execute("INSERT OR IGNORE INTO ppi(t1,t2,score)VALUES(?,?,?)",(a,b,sc))
    for dn,gs in DISEASE_MAP.items():
        c.execute("INSERT OR IGNORE INTO diseases(name)VALUES(?)",(dn,))
        c.execute("SELECT id FROM diseases WHERE name=?",(dn,));did=c.fetchone()[0]
        for g in gs:
            tid=tm.get(g)
            if tid:c.execute("INSERT OR IGNORE INTO target_diseases(target_id,disease_id)VALUES(?,?)",(tid,did))
    conn.commit()
    for aid,aname in c.execute("SELECT id,name FROM ingredients").fetchall():
        for g in IT_MAP.get(aname,[]):
            tid=tm.get(g)
            if tid:c.execute("INSERT OR IGNORE INTO ingredient_targets(ingredient_id,target_id)VALUES(?,?)",(aid,tid))
    conn.commit();conn.close()
    
    # 从 ChEMBL 导入靶点
    chembl_db = os.path.join(BASE, "chembl_36.db")
    if os.path.exists(chembl_db):
        try:
            cb = sqlite3.connect(chembl_db)
            cc = cb.cursor()
            cc.execute("SELECT DISTINCT td.chembl_id, td.pref_name FROM target_dictionary td JOIN assays ass ON td.tid=ass.tid JOIN activities act ON ass.assay_id=act.assay_id WHERE td.organism='Homo sapiens' LIMIT 5000")
            conn2 = get_db()
            c2 = conn2.cursor()
            added = 0
            for cid, pname in cc.fetchall():
                if not pname or len(pname)>80: continue
                gene = pname.split('-')[0].split('/')[0].strip()[:12].replace(' ','_').replace("'","").replace('.','')
                c2.execute("INSERT OR IGNORE INTO targets(gene_name,uniprot_id,protein_name,pubmed_total,pubmed_disease,novelty_score,has_pocket,pocket_score) VALUES(?,?,?,100,50,0.5,0,0)", (gene, cid, pname[:80]))
                if c2.lastrowid: added += 1
            conn2.commit();conn2.close()
            cb.close()
            print(f"  ChEMBL targets: {added} imported")
        except Exception as e:
            print(f"  ChEMBL targets: skipped ({e})")

def _pubmed_query(query):
    time.sleep(0.4)
    url="https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi?"
    if os.environ.get("NCBI_API_KEY"):
        url+="&api_key="+os.environ["NCBI_API_KEY"]
    params=urllib.parse.urlencode({"db":"pubmed","term":query,"rettype":"count","retmode":"xml"})
    try:
        with urllib.request.urlopen(url+params,timeout=10)as resp:
            root=ET.fromstring(resp.read());el=root.find(".//Count")
            return int(el.text)if el is not None else 0
    except:return -1

def get_pubmed_counts(gene_name):
    conn=sqlite3.connect(DB_PATH);c=conn.cursor()
    c.execute("CREATE TABLE IF NOT EXISTS pubmed_cache(gene TEXT PRIMARY KEY,total INTEGER,disease INTEGER,updated INTEGER)")
    now=int(time.time());c.execute("SELECT total,disease,updated FROM pubmed_cache WHERE gene=?",(gene_name,));row=c.fetchone()
    if row and(now-row[2])<7*24*3600:conn.close();return row[0],row[1]
    total=_pubmed_query(gene_name+"[sym]")
    disease=_pubmed_query(gene_name+"[sym] AND(disease OR cancer OR tumor OR diabetes OR Alzheimer OR inflammation OR carcinoma)")
    if total>=0 and disease>=0:c.execute("INSERT OR REPLACE INTO pubmed_cache VALUES(?,?,?,?)",(gene_name,total,disease,now));conn.commit()
    conn.close();return total,disease

def pipeline(input_names):
    # v2.1: herb_id filtered
    conn=get_db();c=conn.cursor();aids=set();herb_ids=set()
    for name in input_names:
        c.execute("SELECT id FROM herbs WHERE chinese_name=? OR latin_name=?",(name,name));hr=c.fetchone()
        if not hr:continue
        herb_ids.add(hr[0])
        c.execute("SELECT id FROM ingredients WHERE herb_id=?",(hr[0],))
        for r in c.fetchall():aids.add(r[0])
    if not aids:conn.close();return{"target_count":0,"novel_target_count":0,"targets":[]}
    ph=",".join(["?"]*len(aids))
    sql="SELECT DISTINCT t.id,t.gene_name,t.uniprot_id,t.protein_name,t.pubmed_total,t.pubmed_disease,t.novelty_score,t.has_pocket,t.pocket_score FROM ingredient_targets it JOIN targets t ON it.target_id=t.id WHERE it.ingredient_id IN("+ph+")"
    c.execute(sql,list(aids));rows=c.fetchall();targets=[dict(r)for r in rows]
    if not targets:conn.close();return{"target_count":0,"novel_target_count":0,"targets":[]}
    for t in targets:
        gene=t["gene_name"]
        total,disease=get_pubmed_counts(gene)
        if total>0:
            t["pubmed_total"]=total;t["pubmed_disease"]=disease
            lf=1.0-math.log10(total)/4.0
            if lf<0:lf=0.0
            if total>0:df=1.0-disease/total
            else:df=1.0
            if df<0:df=0.0
            t["novelty_score"]=round(0.4*lf+0.6*df,4)
        else:
            pc=PUBMED_COUNTS.get(gene,(500,200))
            tt=pc[0]
            if tt<1:tt=1
            dd=pc[1]
            lf=1.0-math.log10(tt)/4.0
            if lf<0:lf=0.0
            if tt>0:df=1.0-dd/tt
            else:df=1.0
            if df<0:df=0.0
            t["novelty_score"]=round(0.4*lf+0.6*df,4)
    tids=[t["id"]for t in targets];tgs={t["id"]:t["gene_name"]for t in targets}
    pht=",".join(["?"]*len(tids))
    c.execute("SELECT t1,t2 FROM ppi WHERE t1 IN("+pht+")OR t2 IN("+pht+")",tids+tids);pr=c.fetchall()
    c.execute("SELECT td.target_id,d.name FROM target_diseases td JOIN diseases d ON td.disease_id=d.id WHERE td.target_id IN("+pht+")",tids);dr=c.fetchall()
    nei={}
    for r in pr:
        a,b=r["t1"],r["t2"]
        if a in tgs:nei.setdefault(a,set()).add(b)
        if b in tgs:nei.setdefault(b,set()).add(a)
    dtg={}
    for r in dr:dtg.setdefault(r["target_id"],[]).append(r["name"])
    for t in targets:
        nbs=nei.get(t["id"],set());np_=sum(1 for n in nbs if n in dtg)
        t["ppi_score"]=round(np_/max(len(nbs),1),4)if nbs else 0;t["disease_partner_count"]=np_
        hits={}
        for n in nbs:
            for d in dtg.get(n,[]):hits[d]=hits.get(d,0)+1
        t["diseases"]=sorted(hits.items(),key=lambda x:x[1],reverse=True)[:5]
    for t in targets:
        t["binding"]=None  # 不估算，需实际对接
    # 按新颖性排序（不做加权合成）
    targets.sort(key=lambda x:x["novelty_score"],reverse=True)
    for t in targets:
        if herb_ids:
            _hph1=",".join(["?"]*len(herb_ids))
            c.execute("SELECT COUNT(*)FROM ingredient_targets it JOIN ingredients i ON it.ingredient_id=i.id WHERE it.target_id=? AND i.herb_id IN("+_hph1+")",[t["id"]]+list(herb_ids));ic=c.fetchone()[0]
        else:
            c.execute("SELECT COUNT(*)FROM ingredient_targets WHERE target_id=?",(t["id"],));ic=c.fetchone()[0]
        pe=1 if t["has_pocket"]else 0;pie=min(t.get("disease_partner_count",0),5)/5.0
        rl=(min(ic,5)/5.0)*0.4+pe*0.3+pie*0.3;t["reliability_score"]=round(rl,2);t["ingredient_count"]=ic
    for t in targets:
        if herb_ids:
            _hph2=",".join(["?"]*len(herb_ids))
            c.execute("SELECT i.name,i.smiles,i.passes_adme FROM ingredient_targets it JOIN ingredients i ON it.ingredient_id=i.id WHERE it.target_id=? AND i.herb_id IN("+_hph2+") ORDER BY i.passes_adme DESC, it.confidence DESC",[t["id"]]+list(herb_ids))
        else:
            c.execute("SELECT i.name,i.smiles,i.passes_adme FROM ingredient_targets it JOIN ingredients i ON it.ingredient_id=i.id WHERE it.target_id=? ORDER BY i.passes_adme DESC, it.confidence DESC",(t["id"],))
        rows = c.fetchall()
        t["_rows"] = rows
        seen = set()
        t["ingredients"]=[]
        for r in rows:
            n=r[0]+(" \u2713"if r[2]else" \u2717")
            if n not in seen:
                seen.add(n)
                t["ingredients"].append(n)
        t["top_smiles"] = rows[0][1] if rows and rows[0][1] else ""
        t["top_ingredient"] = rows[0][0] if rows else ""
    ing_smiles = {}
    for t in targets:
        for r in t.get("_rows", []):
            if r[1]: ing_smiles[r[0]] = r[1]
        del t["_rows"]  # Row 对象不可 JSON 序列化
    conn.close();nc=sum(1 for t in targets if t["novelty_score"]>=0.6)
    return{"target_count":len(targets),"novel_target_count":nc,"targets":targets,"ingredient_smiles":ing_smiles}

@app.route("/")
def index():return serve_html()
@app.route("/api/v1/health")
def health():return{"status":"ok","version":"1.0","name":"TCM-TargetMiner"}
@app.route("/api/v1/herbs")
def list_herbs():
    conn=get_db();c=conn.cursor();c.execute("SELECT id,chinese_name,latin_name,category FROM herbs ORDER BY chinese_name");result=[]
    for r in c.fetchall():
        c.execute("SELECT COUNT(*)FROM ingredients WHERE herb_id=?",(r["id"],));total=c.fetchone()[0]
        c.execute("SELECT COUNT(*)FROM ingredients WHERE herb_id=? AND passes_adme=1",(r["id"],));active=c.fetchone()[0]
        result.append({"id":r["id"],"name":r["chinese_name"],"latin":r["latin_name"],"category":r["category"],"total":total,"active":active})
    conn.close();return jsonify(result)
@app.route("/api/v1/herbs/search")
def search_herbs():
    q=request.args.get("q","")
    if not q:return jsonify([])
    conn=get_db();c=conn.cursor();like="%"+q+"%"
    c.execute("SELECT id,chinese_name,latin_name,category FROM herbs WHERE chinese_name LIKE ? OR latin_name LIKE ? LIMIT 10",(like,like));result=[]
    for r in c.fetchall():
        c.execute("SELECT COUNT(*)FROM ingredients WHERE herb_id=? AND passes_adme=1",(r["id"],));active=c.fetchone()[0]
        c.execute("SELECT COUNT(*)FROM ingredients WHERE herb_id=?",(r["id"],));total=c.fetchone()[0]
        result.append({"id":r["id"],"chinese_name":r["chinese_name"],"latin_name":r["latin_name"],"category":r["category"],"ingredient_count":total,"active_ingredient_count":active})
    conn.close();return jsonify(result)
FORMULAS_DATA = [
    # 解表 (10)
    ("麻黄汤", ["麻黄","桂枝","苦杏仁","甘草"], "解表", "发汗解表宣肺平喘"),
    ("桂枝汤", ["桂枝","白芍","甘草","生姜","大枣"], "解表", "解肌发表调和营卫"),
    ("小青龙汤", ["麻黄","桂枝","细辛","干姜","白芍","半夏","甘草","五味子"], "解表", "解表散寒温肺化饮"),
    ("九味羌活汤", ["羌活","防风","苍术","细辛","川芎","白芷","生地黄","黄芩","甘草"], "解表", "发汗祛湿兼清里热"),
    ("银翘散", ["金银花","连翘","桔梗","薄荷","牛蒡子","甘草"], "解表", "辛凉透表清热解毒"),
    ("桑菊饮", ["桑叶","菊花","苦杏仁","连翘","薄荷","桔梗","甘草"], "解表", "疏风清热宣肺止咳"),
    ("葛根汤", ["葛根","麻黄","桂枝","白芍","甘草","生姜","大枣"], "解表", "发汗解表升津舒经"),
    ("柴葛解肌汤", ["柴胡","葛根","黄芩","白芍","桔梗","甘草","生姜","大枣"], "解表", "解肌清热"),
    ("升麻葛根汤", ["升麻","葛根","白芍","甘草"], "解表", "解肌透疹"),
    ("麻黄附子细辛汤", ["麻黄","细辛"], "解表", "助阳解表"),
    # 泻下 (6)
    ("大承气汤", ["大黄","厚朴","枳实"], "泻下", "峻下热结"),
    ("大黄牡丹汤", ["大黄","牡丹皮","桃仁","芒硝"], "泻下", "泻热破瘀散结消肿"),
    ("温脾汤", ["大黄","人参","甘草","干姜"], "泻下", "温补脾阳攻下冷积"),
    ("麻子仁丸", ["大黄","枳实","厚朴","白芍","苦杏仁"], "泻下", "润肠泄热行气通便"),
    ("济川煎", ["当归","牛膝","枳实","泽泻"], "泻下", "温肾益精润肠通便"),
    ("大黄甘草汤", ["大黄","甘草"], "泻下", "泻热去实"),
    # 和解 (7)
    ("小柴胡汤", ["柴胡","黄芩","人参","半夏","甘草","生姜","大枣"], "和解", "和解少阳"),
    ("大柴胡汤", ["柴胡","黄芩","白芍","半夏","枳实","大黄","生姜","大枣"], "和解", "和解少阳内泻热结"),
    ("逍遥散", ["柴胡","当归","白芍","白术","茯苓","甘草"], "和解", "疏肝解郁养血健脾"),
    ("半夏泻心汤", ["半夏","黄芩","黄连","干姜","人参","甘草","大枣"], "和解", "寒热平调消痞散结"),
    ("痛泻要方", ["白术","白芍","防风","陈皮"], "和解", "补脾柔肝祛湿止泻"),
    ("四逆散", ["柴胡","白芍","枳实","甘草"], "和解", "透邪解郁疏肝理脾"),
    ("达原饮", ["槟榔","厚朴","知母","白芍","黄芩","甘草"], "和解", "开达膜原辟秽化浊"),
    # 清热 (21)
    ("白虎汤", ["石膏","知母","甘草"], "清热", "清热生津"),
    ("黄连解毒汤", ["黄连","黄芩","黄柏","栀子"], "清热", "泻火解毒"),
    ("龙胆泻肝汤", ["龙胆","黄芩","栀子","泽泻","当归","生地黄","柴胡","甘草"], "清热", "清肝胆湿热"),
    ("导赤散", ["生地黄","甘草"], "清热", "清心利水养阴"),
    ("清营汤", ["生地黄","玄参","麦冬","丹参","黄连","金银花","连翘"], "清热", "清营解毒透热养阴"),
    ("犀角地黄汤", ["生地黄","赤芍","牡丹皮"], "清热", "清热解毒凉血散瘀"),
    ("普济消毒饮", ["黄芩","黄连","牛蒡子","连翘","薄荷","桔梗","甘草"], "清热", "清热解毒疏风散邪"),
    ("仙方活命饮", ["金银花","赤芍","白芷","防风","甘草","陈皮"], "清热", "清热解毒消肿溃坚"),
    ("葛根芩连汤", ["葛根","黄芩","黄连","甘草"], "清热", "解表清里"),
    ("芍药汤", ["白芍","当归","黄连","黄芩","大黄","甘草","肉桂","槟榔"], "清热", "清热燥湿调气和血"),
    ("白头翁汤", ["黄柏","黄连"], "清热", "清热解毒凉血止痢"),
    ("左金丸", ["黄连","吴茱萸"], "清热", "清肝泻火降逆止呕"),
    ("泻心汤", ["大黄","黄连","黄芩"], "清热", "泻火解毒燥湿泄痞"),
    ("凉膈散", ["大黄","芒硝","甘草","栀子","薄荷","黄芩","连翘"], "清热", "泻火通便清上泄下"),
    ("泻白散", ["桑白皮","地骨皮","甘草"], "清热", "清泻肺热平喘止咳"),
    ("玉女煎", ["石膏","熟地黄","麦冬","知母","牛膝"], "清热", "清胃热滋肾阴"),
    ("青蒿鳖甲汤", ["青蒿","知母","生地黄","牡丹皮"], "清热", "养阴透热"),
    ("五味消毒饮", ["金银花","蒲公英","紫花地丁"], "清热", "清热解毒消散疔疮"),
    ("四妙勇安汤", ["金银花","玄参","当归","甘草"], "清热", "清热解毒活血止痛"),
    ("桔梗汤", ["桔梗","甘草"], "清热", "宣肺利咽"),
    ("栀子豉汤", ["栀子"], "清热", "清热除烦"),
    # 温里 (9)
    ("理中丸", ["人参","干姜","白术","甘草"], "温里", "温中祛寒补气健脾"),
    ("小建中汤", ["桂枝","白芍","甘草","生姜","大枣"], "温里", "温中补虚和里缓急"),
    ("四逆汤", ["甘草","干姜"], "温里", "回阳救逆"),
    ("当归四逆汤", ["当归","桂枝","白芍","细辛","甘草","大枣"], "温里", "温经散寒养血通脉"),
    ("阳和汤", ["熟地黄","肉桂","麻黄","鹿茸","甘草"], "温里", "温阳补血散寒通滞"),
    ("吴茱萸汤", ["吴茱萸","人参","生姜","大枣"], "温里", "温中补虚降逆止呕"),
    ("大建中汤", ["花椒","干姜","人参"], "温里", "温中补虚降逆止痛"),
    ("黄芪桂枝五物汤", ["黄芪","桂枝","白芍","生姜","大枣"], "温里", "益气温经和血通痹"),
    ("甘草干姜汤", ["甘草","干姜"], "温里", "温中益气"),
    # 补益 (19)
    ("四君子汤", ["人参","白术","茯苓","甘草"], "补益", "益气健脾"),
    ("六君子汤", ["人参","白术","茯苓","甘草","陈皮","半夏"], "补益", "益气健脾燥湿化痰"),
    ("补中益气汤", ["黄芪","人参","白术","甘草","当归","陈皮","柴胡"], "补益", "补中益气升阳举陷"),
    ("生脉散", ["人参","麦冬","五味子"], "补益", "益气生津敛阴止汗"),
    ("玉屏风散", ["黄芪","白术","防风"], "补益", "益气固表止汗"),
    ("四物汤", ["当归","川芎","白芍","熟地黄"], "补益", "补血调血"),
    ("当归补血汤", ["黄芪","当归"], "补益", "补气生血"),
    ("归脾汤", ["人参","黄芪","白术","茯苓","甘草","当归","酸枣仁","远志"], "补益", "益气补血健脾养心"),
    ("八珍汤", ["人参","白术","茯苓","甘草","当归","白芍","川芎","熟地黄"], "补益", "益气补血"),
    ("炙甘草汤", ["甘草","生姜","人参","生地黄","桂枝","阿胶","麦冬","大枣"], "补益", "益气滋阴通阳复脉"),
    ("六味地黄丸", ["熟地黄","山茱萸","山药","牡丹皮","茯苓"], "补益", "滋补肾阴"),
    ("左归丸", ["熟地黄","山药","山茱萸","枸杞子","牛膝","鹿茸"], "补益", "滋阴补肾填精益髓"),
    ("右归丸", ["熟地黄","山药","山茱萸","枸杞子","杜仲","肉桂","当归"], "补益", "温补肾阳填精益髓"),
    ("一贯煎", ["生地黄","枸杞子","当归","麦冬"], "补益", "滋阴疏肝"),
    ("肾气丸", ["熟地黄","山茱萸","山药","茯苓","牡丹皮","桂枝"], "补益", "补肾助阳"),
    ("参苓白术散", ["人参","白术","茯苓","甘草","山药","莲子","砂仁","桔梗","薏苡仁"], "补益", "益气健脾渗湿止泻"),
    ("大补阴丸", ["熟地黄","龟甲","黄柏","知母"], "补益", "滋阴降火"),
    ("二至丸", ["女贞子","墨旱莲"], "补益", "补益肝肾滋阴止血"),
    ("七宝美髯丹", ["何首乌","茯苓","牛膝","当归","枸杞子","菟丝子","补骨脂"], "补益", "补益肝肾乌发壮骨"),
    # 安神 (4)
    ("酸枣仁汤", ["酸枣仁","知母","川芎","茯苓","甘草"], "安神", "养血安神清热除烦"),
    ("天王补心丹", ["酸枣仁","生地黄","玄参","麦冬","丹参","当归","茯苓","远志","五味子"], "安神", "滋阴养血补心安神"),
    ("甘麦大枣汤", ["甘草","大枣"], "安神", "养心安神和中缓急"),
    ("交泰丸", ["黄连","肉桂"], "安神", "交通心肾"),
    # 理气 (7)
    ("越鞠丸", ["川芎","栀子","白术"], "理气", "行气解郁"),
    ("半夏厚朴汤", ["半夏","厚朴","茯苓","生姜"], "理气", "行气散结降逆化痰"),
    ("苏子降气汤", ["半夏","当归","甘草","肉桂"], "理气", "降气平喘祛痰止咳"),
    ("定喘汤", ["麻黄","苦杏仁","黄芩","半夏","甘草"], "理气", "宣肺降气清热化痰"),
    ("柴胡疏肝散", ["柴胡","陈皮","川芎","香附","枳实","白芍","甘草"], "理气", "疏肝理气活血止痛"),
    ("四磨汤", ["人参","槟榔","乌药"], "理气", "行气降逆宽胸散结"),
    ("厚朴温中汤", ["厚朴","陈皮","甘草","茯苓","木香","干姜"], "理气", "温中行气燥湿除满"),
    # 理血 (8)
    ("桃红四物汤", ["桃仁","红花","当归","白芍","川芎","熟地黄"], "理血", "活血养血"),
    ("血府逐瘀汤", ["当归","红花","柴胡","甘草","川芎","枳实"], "理血", "活血祛瘀行气止痛"),
    ("补阳还五汤", ["黄芪","当归","川芎","红花","桃仁","赤芍"], "理血", "补气活血通络"),
    ("复元活血汤", ["柴胡","天花粉","当归","红花","甘草","大黄","桃仁"], "理血", "活血祛瘀疏肝通络"),
    ("丹参饮", ["丹参"], "理血", "活血祛瘀行气止痛"),
    ("温经汤", ["当归","白芍","川芎","人参","桂枝","甘草"], "理血", "温经散寒养血祛瘀"),
    ("生化汤", ["当归","川芎","桃仁","甘草"], "理血", "活血化瘀温经止痛"),
    ("桂枝茯苓丸", ["桂枝","茯苓","牡丹皮","桃仁","白芍"], "理血", "活血化瘀缓消癥块"),
    # 治风 (3)
    ("川芎茶调散", ["川芎","荆芥","防风","细辛","白芷","薄荷","甘草"], "治风", "疏风止痛"),
    ("天麻钩藤饮", ["天麻","钩藤","牛膝","栀子","黄芩","杜仲"], "治风", "平肝熄风清热活血"),
    ("大秦艽汤", ["秦艽","甘草","川芎","当归","白芍","细辛","羌活","防风","黄芩","石膏","白芷","白术","生地黄","熟地黄","茯苓","独活"], "治风", "祛风清热养血活血"),
    # 治燥 (4)
    ("麦门冬汤", ["麦冬","半夏","人参","甘草","大枣"], "治燥", "滋养肺胃降逆下气"),
    ("养阴清肺汤", ["生地黄","麦冬","玄参","白芍","牡丹皮","薄荷","甘草"], "治燥", "养阴清肺"),
    ("增液汤", ["玄参","麦冬","生地黄"], "治燥", "增液润燥"),
    ("百合固金汤", ["百合","生地黄","熟地黄","麦冬","玄参","当归","白芍","桔梗","甘草","川贝母"], "治燥", "滋养肺肾止咳化痰"),
    # 祛湿 (12)
    ("平胃散", ["苍术","厚朴","陈皮","甘草"], "祛湿", "燥湿运脾行气和胃"),
    ("藿香正气散", ["藿香","紫苏","白芷","半夏","厚朴","陈皮","白术","茯苓","桔梗","甘草","生姜","大枣"], "祛湿", "解表化湿理气和中"),
    ("茵陈蒿汤", ["茵陈"], "祛湿", "清热利湿退黄"),
    ("五苓散", ["猪苓","茯苓","泽泻","白术","桂枝"], "祛湿", "利水渗湿温阳化气"),
    ("苓桂术甘汤", ["茯苓","桂枝","白术","甘草"], "祛湿", "温阳化饮健脾利湿"),
    ("真武汤", ["茯苓","白芍","白术","生姜"], "祛湿", "温阳利水"),
    ("实脾散", ["厚朴","白术","木瓜","木香","槟榔","茯苓","干姜","甘草"], "祛湿", "温阳健脾行气利水"),
    ("二妙散", ["黄柏","苍术"], "祛湿", "清热燥湿"),
    ("羌活胜湿汤", ["羌活","独活","防风","川芎","甘草"], "祛湿", "祛风胜湿止痛"),
    ("防己黄芪汤", ["防己","黄芪","白术","甘草","生姜","大枣"], "祛湿", "益气祛风健脾利水"),
    ("鸡鸣散", ["槟榔","木瓜","吴茱萸","桔梗","生姜","紫苏"], "祛湿", "行气降浊宣化寒湿"),
    ("连朴饮", ["黄连","厚朴","半夏","栀子"], "祛湿", "清热化湿理气和中"),
    # 祛痰 (6)
    ("二陈汤", ["半夏","陈皮","茯苓","甘草"], "祛痰", "燥湿化痰理气和中"),
    ("温胆汤", ["半夏","竹茹","枳实","陈皮","茯苓","甘草"], "祛痰", "理气化痰清胆和胃"),
    ("清气化痰丸", ["黄芩","半夏","苦杏仁","枳实","陈皮","茯苓"], "祛痰", "清热化痰理气止咳"),
    ("苓甘五味姜辛汤", ["茯苓","甘草","干姜","细辛","五味子"], "祛痰", "温肺化饮"),
    ("半夏白术天麻汤", ["半夏","天麻","白术","茯苓","陈皮","甘草","生姜","大枣"], "祛痰", "化痰熄风健脾祛湿"),
    ("止嗽散", ["桔梗","荆芥","紫菀","百部","白前","甘草","陈皮"], "祛痰", "宣利肺气疏风止咳"),
    # 消食 (3)
    ("保和丸", ["山楂","神曲","半夏","茯苓","陈皮","莱菔子"], "消食", "消食和胃"),
    ("健脾丸", ["白术","茯苓","人参","山药","山楂","神曲","麦芽","木香","砂仁","陈皮","黄连","甘草"], "消食", "健脾和胃消食止泻"),
    ("枳术丸", ["枳实","白术"], "消食", "健脾消痞"),
    # 驱虫 (1)
    ("乌梅丸", ["乌梅","细辛","黄连","当归","花椒","桂枝","人参","黄柏","干姜"], "驱虫", "温脏安蛔"),
    # 止痛 (1)
    ("芍药甘草汤", ["白芍","甘草"], "止痛", "缓急止痛"),
]

@app.route("/api/v1/formulas")


def list_formulas():
    result=[]
    for n,h,c,nt in FORMULAS_DATA:result.append({"name":n,"herbs":h,"category":c,"note":nt})
    return jsonify(result)
@app.route("/api/v1/analyze",methods=["POST"])
def start_analysis():
    data=request.get_json()or{};names=data.get("names")or data.get("input_names")or[]
    if not names:return jsonify({"error":"Need at least one herb name"}),400
    tid=str(uuid.uuid4())[:8];result=pipeline(names)
    conn=get_db();c=conn.cursor()
    c.execute("INSERT OR REPLACE INTO results(task_id,input_type,input_name,result_json)VALUES(?,?,?,?)",(tid,"herb","+".join(names),json.dumps(result,ensure_ascii=False)))
    conn.commit();conn.close();result["task_id"]=tid;result["status"]="success";return jsonify(result)
@app.route("/api/v1/task/<tid>")
def task_status(tid):
    conn=get_db();c=conn.cursor();c.execute("SELECT result_json FROM results WHERE task_id=?",(tid,));row=c.fetchone();conn.close()
    if not row:return jsonify({"status":"not_found"})
    rj=json.loads(row[0]);return jsonify({"task_id":tid,"status":"success","target_count":rj.get("target_count",0),"novel_target_count":rj.get("novel_target_count",0)})
@app.route("/api/v1/task/<tid>/result")
def task_result(tid):
    conn=get_db();c=conn.cursor();c.execute("SELECT*FROM results WHERE task_id=?",(tid,));row=c.fetchone();conn.close()
    if not row:return jsonify({"error":"not found"}),404
    rj=json.loads(row["result_json"]);return jsonify({"task_id":row["task_id"],"status":"success","target_count":rj.get("target_count",0),"novel_target_count":rj.get("novel_target_count",0),"targets":rj.get("targets",[])})
@app.route("/api/v1/predict_molecule",methods=["POST"])
def predict_molecule():
    data=request.get_json()or{};smiles=data.get("smiles","").strip()
    if not smiles:return jsonify({"error":"Need SMILES string"}),400
    try:
        result=predict_molecule_targets(smiles)
        if"error"in result:return jsonify(result),400
        ppi=infer_ppi_targets({t['gene_name']:t['confidence']for t in result['targets']},result['similar_compounds'])
        if ppi:result['targets'].extend(ppi);result['target_count']=len(result['targets'])
        return jsonify(result)
    except Exception as e:return jsonify({"error":"Prediction failed: "+str(e)}),500
@app.route("/api/v1/molecule_from_smiles",methods=["POST"])
def molecule_from_smiles():
    data=request.get_json()or{};smiles=data.get("smiles","").strip()
    if not smiles:return jsonify({"error":"Need SMILES string"}),400
    fp=ecfp_fingerprint(smiles)
    if fp is None:return jsonify({"error":"Cannot parse SMILES"}),400
    return jsonify({"status":"ok"})

if __name__=="__main__":
    # 清理临时文件
    import glob as _g
    for f in _g.glob(os.path.join(BASE,"_d_*")):
        try: os.remove(f)
        except: pass
    init_db()
    # 导入 ChEMBL 数据
    chembl_file = os.path.join(BASE, "chembl_itmap.json")
    if os.path.exists(chembl_file) and not os.path.exists(os.path.join(BASE, "chembl_imported.flag")):
        print("  Importing ChEMBL data...")
        with open(chembl_file, "r", encoding="utf-8") as f:
            chembl_data = json.load(f)
        conn = get_db()
        c = conn.cursor()
        c.execute("SELECT id, chinese_name FROM herbs WHERE chinese_name='参考库'")
        ref = c.fetchone()
        ref_id = ref['id'] if ref else 9999  # 避免落入半枝莲(id=1)
        c.execute("SELECT id, gene_name FROM targets")
        tgt_map = {r['gene_name']: r['id'] for r in c.fetchall()}
        added_compounds = 0
        added_targets = 0
        for key, data in chembl_data.items():
            smi = data.get("smiles", "")
            if not smi or len(smi) > 300:
                continue
            name = (data.get("name") or f"ChEMBL_{added_compounds:06d}")[:80]
            c.execute("INSERT OR IGNORE INTO ingredients (herb_id, name, smiles, ob, dl, mw, logp, passes_adme) VALUES (?,?,?,30.0,0.18,300.0,2.0,1)", (ref_id, name, smi))
            if c.lastrowid:
                added_compounds += 1
                iid = c.lastrowid
                for tname in data.get("targets", [])[:5]:
                    # 尝试匹配已有靶点
                    for g, tid in tgt_map.items():
                        if g.lower() in tname.lower() or tname.lower() in g.lower():
                            c.execute("INSERT OR IGNORE INTO ingredient_targets (ingredient_id, target_id) VALUES (?,?)", (iid, tid))
                            added_targets += 1
                            break
            if added_compounds % 1000 == 0:
                conn.commit()
        conn.commit()
        conn.close()
        with open(os.path.join(BASE, "chembl_imported.flag"), "w") as f:
            f.write("done")
        print(f"  ChEMBL import: {added_compounds} compounds, {added_targets} target links")
    auto_expand_itmap()
    ppi_c=sum(len(v)for v in PPI_MAP.values())
    print("="*55)
    print("  TCM-TargetMiner v1.0")
    print("  Herbs: %d(TCM+Miao)"%len(HERBS_DATA))
    print("  Targets: "+str(len(TARGETS_DATA)))
    print("  PPI edges: "+str(ppi_c))
    print("="*55)
    print("  >>> http://localhost:5000")
    print("="*55)
    app.run(host="0.0.0.0",port=5000,debug=True)



def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    return conn

def init_db():
    conn = get_db()
    c = conn.cursor()
    c.executescript("""
        CREATE TABLE IF NOT EXISTS herbs (id INTEGER PRIMARY KEY AUTOINCREMENT, chinese_name TEXT UNIQUE NOT NULL, latin_name TEXT, category TEXT, properties TEXT);
        CREATE TABLE IF NOT EXISTS ingredients (id INTEGER PRIMARY KEY AUTOINCREMENT, herb_id INTEGER NOT NULL, name TEXT NOT NULL, smiles TEXT, ob REAL, dl REAL, mw REAL, logp REAL, passes_adme INTEGER DEFAULT 1, FOREIGN KEY (herb_id) REFERENCES herbs(id));
        CREATE TABLE IF NOT EXISTS targets (id INTEGER PRIMARY KEY AUTOINCREMENT, gene_name TEXT UNIQUE NOT NULL, uniprot_id TEXT, protein_name TEXT, pubmed_total INTEGER DEFAULT 0, pubmed_disease INTEGER DEFAULT 0, novelty_score REAL DEFAULT 0.5, has_pocket INTEGER, pocket_score REAL DEFAULT 0);
        CREATE TABLE IF NOT EXISTS ingredient_targets (id INTEGER PRIMARY KEY AUTOINCREMENT, ingredient_id INTEGER NOT NULL, target_id INTEGER NOT NULL, confidence REAL DEFAULT 1.0, UNIQUE(ingredient_id, target_id), FOREIGN KEY (ingredient_id) REFERENCES ingredients(id), FOREIGN KEY (target_id) REFERENCES targets(id));
        CREATE TABLE IF NOT EXISTS ppi (id INTEGER PRIMARY KEY AUTOINCREMENT, t1 INTEGER NOT NULL, t2 INTEGER NOT NULL, score INTEGER DEFAULT 900, UNIQUE(t1, t2), FOREIGN KEY (t1) REFERENCES targets(id), FOREIGN KEY (t2) REFERENCES targets(id));
        CREATE TABLE IF NOT EXISTS diseases (id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT UNIQUE NOT NULL);
        CREATE TABLE IF NOT EXISTS target_diseases (id INTEGER PRIMARY KEY AUTOINCREMENT, target_id INTEGER NOT NULL, disease_id INTEGER NOT NULL, UNIQUE(target_id, disease_id), FOREIGN KEY (target_id) REFERENCES targets(id), FOREIGN KEY (disease_id) REFERENCES diseases(id));
        CREATE TABLE IF NOT EXISTS results (id INTEGER PRIMARY KEY AUTOINCREMENT, task_id TEXT UNIQUE NOT NULL, input_type TEXT, input_name TEXT, result_json TEXT, created_at TEXT DEFAULT (datetime('now')));
    """)
    c.execute("SELECT COUNT(*) FROM herbs")
    if c.fetchone()[0] > 0:
        conn.close()
        return
    for h in HERBS_DATA:
        c.execute("INSERT INTO herbs (chinese_name, latin_name, category, properties) VALUES (?,?,?,?)", h)
    i = 0
    for gene, (up, pname) in TARGETS_DATA.items():
        pc = PUBMED_COUNTS.get(gene, (500, 200))
        t = pc[0]
        if t < 1:
            t = 1
        d = pc[1]
        lf = 1.0 - math.log10(t) / 4.0
        if lf < 0:
            lf = 0.0
        if t > 0:
            df = 1.0 - d / t
        else:
            df = 1.0
        if df < 0:
            df = 0.0
        ns = round(0.4 * lf + 0.6 * df, 4)
        pk = POCKET_DATA.get(gene, (0, 0.0))
        c.execute("INSERT INTO targets (gene_name, uniprot_id, protein_name, pubmed_total, pubmed_disease, novelty_score, has_pocket, pocket_score) VALUES (?,?,?,?,?,?,?,?)", (gene, up, pname, pc[0], pc[1], ns, pk[0], pk[1]))
        i += 1
        if i % 10 == 0: conn.commit()
    c.execute("SELECT id, chinese_name FROM herbs")
    herb_ids = {r[1]: r[0] for r in c.fetchall()}
    for row in INGREDIENTS_DATA:
        h_id = herb_ids.get(row[0])
        if not h_id: continue
        passes = 1 if (row[3] >= 30 and row[4] >= 0.18) else 0
        c.execute("INSERT INTO ingredients (herb_id, name, smiles, ob, dl, mw, logp, passes_adme) VALUES (?,?,?,?,?,?,?,?)", (h_id, row[1], row[2], row[3], row[4], row[5], row[6], passes))
    c.execute("SELECT id, name FROM ingredients")
    ing_map = {r[1]: r[0] for r in c.fetchall()}
    c.execute("SELECT id, gene_name FROM targets")
    tgt_map = {r[1]: r[0] for r in c.fetchall()}
    for ing_name, ing_id in ing_map.items():
        genes = IT_MAP.get(ing_name, [])
        for g in genes:
            tid = tgt_map.get(g)
            if tid: c.execute("INSERT OR IGNORE INTO ingredient_targets (ingredient_id, target_id) VALUES (?,?)", (ing_id, tid))
    for g1, partners in PPI_MAP.items():
        tid1 = tgt_map.get(g1)
        if not tid1: continue
        for g2 in partners:
            tid2 = tgt_map.get(g2)
            if not tid2: continue
            a, b = min(tid1, tid2), max(tid1, tid2)
            c.execute("INSERT OR IGNORE INTO ppi (t1, t2) VALUES (?,?)", (a, b))
    for dname, genes in DISEASE_MAP.items():
        c.execute("INSERT OR IGNORE INTO diseases (name) VALUES (?)", (dname,))
        c.execute("SELECT id FROM diseases WHERE name=?", (dname,))
        did = c.fetchone()[0]
        for g in genes:
            tid = tgt_map.get(g)
            if tid: c.execute("INSERT OR IGNORE INTO target_diseases (target_id, disease_id) VALUES (?,?)", (tid, did))
    conn.commit()
    c.execute("SELECT id, name FROM ingredients")
    for aid, aname in c.fetchall():
        glist = IT_MAP.get(aname, [])
        if not glist: continue
        for g in glist:
            tid = tgt_map.get(g)
            if tid: c.execute("INSERT OR IGNORE INTO ingredient_targets (ingredient_id, target_id) VALUES (?,?)", (aid, tid))
    conn.commit()
    conn.close()

def _pubmed_query(query):
    time.sleep(0.4)
    url = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi?"
    if os.environ.get("NCBI_API_KEY"):
        url += "&api_key=" + os.environ["NCBI_API_KEY"]
    params = urllib.parse.urlencode({"db":"pubmed","term":query,"rettype":"count","retmode":"xml"})
    try:
        with urllib.request.urlopen(url + params, timeout=10) as resp:
            root = ET.fromstring(resp.read())
            el = root.find(".//Count")
            return int(el.text) if el is not None else 0
    except Exception:
        return -1

def get_pubmed_counts(gene_name):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("CREATE TABLE IF NOT EXISTS pubmed_cache (gene TEXT PRIMARY KEY, total INTEGER, disease INTEGER, updated INTEGER)")
    now = int(time.time())
    c.execute("SELECT total, disease, updated FROM pubmed_cache WHERE gene=?", (gene_name,))
    row = c.fetchone()
    if row and (now - row[2]) < 7 * 24 * 3600:
        conn.close()
        return row[0], row[1]
    total = _pubmed_query(gene_name)
    disease = _pubmed_query(gene_name + " AND (disease OR cancer OR tumor)")
    if total >= 0 and disease >= 0:
        c.execute("INSERT OR REPLACE INTO pubmed_cache VALUES (?,?,?,?)", (gene_name, total, disease, now))
        conn.commit()
    conn.close()
    return total, disease

def pipeline(input_names):
    conn = get_db()
    c = conn.cursor()
    all_ing_ids = set()
    herb_ids = set()
    for name in input_names:
        c.execute("SELECT id FROM herbs WHERE chinese_name=? OR latin_name=?", (name, name))
        hr = c.fetchone()
        if not hr: continue
        herb_ids.add(hr[0])
        c.execute("SELECT id FROM ingredients WHERE herb_id=?", (hr[0],))
        for r in c.fetchall(): all_ing_ids.add(r[0])
    if not all_ing_ids:
        conn.close()
        return {"target_count": 0, "novel_target_count": 0, "targets": []}
    ph = ",".join(["?"] * len(all_ing_ids))
    sql = "SELECT DISTINCT t.id, t.gene_name, t.uniprot_id, t.protein_name, t.pubmed_total, t.pubmed_disease, t.novelty_score, t.has_pocket, t.pocket_score FROM ingredient_targets it JOIN targets t ON it.target_id = t.id WHERE it.ingredient_id IN (" + ph + ")"
    c.execute(sql, list(all_ing_ids))
    rows = c.fetchall()
    targets = [dict(r) for r in rows]
    if not targets:
        conn.close()
        return {"target_count": 0, "novel_target_count": 0, "targets": []}
    for t in targets:
        gene = t["gene_name"]
        total, disease = get_pubmed_counts(gene)
        # 兜底：PubMed查不到时用预设值
        if total <= 0:
            pc = PUBMED_COUNTS.get(gene, (500, 200))
            total = pc[0]; disease = pc[1]
        t["pubmed_total"] = total; t["pubmed_disease"] = disease
        if total == 0: t["novelty_score"] = 1.0
        else:
            lf = max(0, 1 - math.log10(max(total, 1)) / 4)
            df_ = max(0, 1 - disease / total) if total > 0 else 1.0
            t["novelty_score"] = round(0.4 * lf + 0.6 * df_, 4)
    tgt_ids = [t["id"] for t in targets]
    tgt_genes = {t["id"]: t["gene_name"] for t in targets}
    pht = ",".join(["?"] * len(tgt_ids))
    c.execute("SELECT t1, t2 FROM ppi WHERE t1 IN (" + pht + ") OR t2 IN (" + pht + ")", tgt_ids + tgt_ids)
    ppi_rows = c.fetchall()
    c.execute("SELECT td.target_id, d.name FROM target_diseases td JOIN diseases d ON td.disease_id=d.id WHERE td.target_id IN (" + pht + ")", tgt_ids)
    dis_rows = c.fetchall()
    nei = {}
    for r in ppi_rows:
        a, b = r["t1"], r["t2"]
        if a in tgt_genes: nei.setdefault(a, set()).add(b)
        if b in tgt_genes: nei.setdefault(b, set()).add(a)
    dis_tgt = {}
    for r in dis_rows: dis_tgt.setdefault(r["target_id"], []).append(r["name"])
    for t in targets:
        nbs = nei.get(t["id"], set())
        np_ = sum(1 for n in nbs if n in dis_tgt)
        t["ppi_score"] = round(np_ / max(len(nbs), 1), 4) if nbs else 0
        t["disease_partner_count"] = np_
        hits = {}
        for n in nbs:
            for d in dis_tgt.get(n, []): hits[d] = hits.get(d, 0) + 1
        t["diseases"] = sorted(hits.items(), key=lambda x: x[1], reverse=True)[:5]
    for t in targets:
        if t["has_pocket"] and t["pocket_score"] > 0.5: t["binding"] = round(-6.0 - t["pocket_score"] * 4.0, 1)
        elif t["has_pocket"]: t["binding"] = round(-4.0 - t["pocket_score"] * 3.0, 1)
        else: t["binding"] = None
    # 按新颖性排序（不做加权合成）
    targets.sort(key=lambda x: x["novelty_score"], reverse=True)
    for t in targets:
        if herb_ids:
            herb_ph2 = ",".join(["?"] * len(herb_ids))
            c.execute("SELECT COUNT(*) FROM ingredient_targets it JOIN ingredients i ON it.ingredient_id=i.id WHERE it.target_id=? AND i.herb_id IN (" + herb_ph2 + ")", [t["id"]] + list(herb_ids))
        else:
            c.execute("SELECT COUNT(*) FROM ingredient_targets WHERE target_id=?", (t["id"],))
        ing_cnt = c.fetchone()[0]
        pocket_ev = min((t.get("pocket_score") or 0) / 0.5, 1.0) if (t["has_pocket"] and (t.get("pocket_score") or -1) >= 0) else 0
        ppi_ev = min(t.get("disease_partner_count", 0), 5) / 5.0
        rel = (min(ing_cnt, 5) / 5.0) * 0.4 + pocket_ev * 0.3 + ppi_ev * 0.3
        if rel >= 0.6: t["reliability"] = "高"
        elif rel >= 0.3: t["reliability"] = "中"
        else: t["reliability"] = "低"
        t["reliability_score"] = round(rel, 2)
        t["ingredient_count"] = ing_cnt
    for t in targets:
        if herb_ids:
            herb_ph3 = ",".join(["?"] * len(herb_ids))
            c.execute("SELECT i.name, i.smiles, i.passes_adme FROM ingredient_targets it JOIN ingredients i ON it.ingredient_id=i.id WHERE it.target_id=? AND i.herb_id IN (" + herb_ph3 + ") ORDER BY i.passes_adme DESC, it.confidence DESC", [t["id"]] + list(herb_ids))
        else:
            c.execute("SELECT i.name, i.smiles, i.passes_adme FROM ingredient_targets it JOIN ingredients i ON it.ingredient_id=i.id WHERE it.target_id=? ORDER BY i.passes_adme DESC, it.confidence DESC", (t["id"],))
        rows = c.fetchall()
        t["_rows"] = rows
        seen3 = set()
        t["ingredients"] = []
        for r in rows:
            n = r[0] + (" ✓" if r[2] else " ✗")
            if n not in seen3:
                seen3.add(n)
                t["ingredients"].append(n)
        t["top_smiles"] = rows[0][1] if rows and rows[0][1] else ""
        t["top_ingredient"] = rows[0][0] if rows else ""
    conn.close()
    nc = sum(1 for t in targets if t["novelty_score"] >= 0.6)
    # 收集所有成分的SMILES
    ing_smiles = {}
    for t in targets:
        rows_all = t.get("_rows", [])
        for r in rows_all:
            ing_smiles[r[0]] = r[1] if r[1] else ""
        del t["_rows"]  # Row 对象不可 JSON 序列化
    return {"target_count": len(targets), "novel_target_count": nc, "targets": targets, "ingredient_smiles": ing_smiles}

from uniprot_pdb import UNIPROT_PDB


@app.route("/api/v1/docking_info/<uniprot>")
def docking_info(uniprot):
    """返回对接所需信息"""
    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT * FROM targets WHERE uniprot_id=?", (uniprot,))
    row = c.fetchone()
    conn.close()
    if not row:
        return jsonify({"error": "Target not found"}), 404
    
    # 找 obabel 和 vina
    vina_path = None
    for p in [os.path.join(BASE,"vina.exe"), r"D:\autodock\vina.exe"]:
        if os.path.exists(p): vina_path = p; break
    
    pdb_url = f"https://alphafold.ebi.ac.uk/files/AF-{uniprot}-F1-model_v4.pdb"
    
    return jsonify({
        "gene_name": row["gene_name"],
        "uniprot": uniprot,
        "pocket_score": row["pocket_score"],
        "has_pocket": bool(row["has_pocket"]),
        "pdb_download": pdb_url,
        "vina_ready": vina_path is not None,
        "docking_cmd": f"vina --receptor AF-{uniprot}-F1-model_v4.pdb --ligand ligand.pdbqt --center_x 0 --center_y 0 --center_z 0 --size_x 20 --size_y 20 --size_z 20" if vina_path else "",
    })


@app.route("/api/pptx_simple2")
def pptx_simple2():
    """Minimal test"""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x27)
    tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = 'Test OK'
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(0x00, 0xBB, 0xF9)
    from pptx.util import Pt
    import io
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Response(buf.read(), mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    headers={'Content-Disposition': 'attachment; filename="test.pptx"'})

@app.route("/api/pptx_direct2")
def pptx_direct2():
    """Generate full 智鉴药靶 PPTX"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE, 'templates'))
        from PptxGen import generate_bytes
        data = generate_bytes()
        return Response(data, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                        headers={'Content-Disposition': 'attachment; filename="智鉴药靶_路演PPT.pptx"'})
    except Exception as e:
        import traceback
        return Response(f"Error: {str(e)}\n{traceback.format_exc()}", status=500)

