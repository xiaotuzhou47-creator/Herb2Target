# -*- coding: utf-8 -*-
"""Generate 智鉴药靶 pitch deck PPTX - v2.0 优化版"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import io

BG_DARK = RGBColor(0x0A, 0x0E, 0x27)
BG_CARD = RGBColor(0x14, 0x1E, 0x38)
BG_HIGHLIGHT = RGBColor(0x00, 0x3D, 0x5C)
C_BLUE = RGBColor(0x00, 0xBB, 0xF9)
C_GREEN = RGBColor(0x00, 0xE6, 0x76)
C_PURPLE = RGBColor(0x7B, 0x2F, 0xFF)
C_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
C_PINK = RGBColor(0xFF, 0x33, 0x66)
C_TEAL = RGBColor(0x00, 0xD4, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xB0, 0xC4, 0xDE)
GRAY = RGBColor(0x66, 0x77, 0x88)
PAGE_GRAY = RGBColor(0x44, 0x55, 0x66)

SLIDE_W = 13.333
SLIDE_H = 7.5

def build_prs():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    sl = prs.slide_layouts[6]

    def dark_bg(s):
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG_DARK

    def tb(s, l, t, w, h, txt, fs=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
        tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tx.text_frame.word_wrap = True
        p = tx.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = b
        p.alignment = a

    def multi_tb(s, l, t, w, h, lines, fs=14, c=LIGHT):
        """Multiple lines in one textbox"""
        tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tx.text_frame.word_wrap = True
        for i, (txt, *opts) in enumerate(lines):
            p = tx.text_frame.paragraphs[0] if i == 0 else tx.text_frame.add_paragraph()
            p.text = txt
            p.font.size = Pt(opts[0] if len(opts) > 0 else fs)
            p.font.color.rgb = opts[1] if len(opts) > 1 else c
            p.font.bold = opts[2] if len(opts) > 2 else False
            if len(opts) > 3:
                p.alignment = opts[3]
            p.space_after = Pt(2)

    def rc(s, l, t, w, h, fc=None, lc=None):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fc or BG_CARD
        if lc:
            sh.line.color.rgb = lc
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    def ac(s, l, t, w, c=None):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Pt(4))
        sh.fill.solid()
        sh.fill.fore_color.rgb = c or C_BLUE
        sh.line.fill.background()

    def page_num(s, n):
        if n <= 1 or n == 12:
            return
        tb(s, SLIDE_W - 1.2, SLIDE_H - 0.5, 1, 0.4, f'{n:02d}', 11, PAGE_GRAY, False, PP_ALIGN.RIGHT)

    def section_title(s, title):
        tb(s, 0.8, 0.35, 8, 0.7, title, 34, WHITE, True)
        ac(s, 0.8, 1.05, 2.2, C_GREEN)

    # ============================================================
    # SLIDE 1: COVER
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    # 副装饰线
    rc(s, 0, SLIDE_H - 0.08, SLIDE_W, 0.08, C_GREEN)
    tb(s, 1, 1.8, 11, 1.5, '智鉴药靶', 56, WHITE, True, PP_ALIGN.CENTER)
    tb(s, 1, 3.2, 11, 0.8, '中药多靶点智能预测与分子对接平台', 26, C_BLUE, False, PP_ALIGN.CENTER)
    ac(s, 5.2, 4.1, 2.933, C_GREEN)
    tb(s, 1, 4.6, 11, 0.6, '贵州中医药大学 · 药学院', 18, GRAY, False, PP_ALIGN.CENTER)
    tb(s, 1, 5.2, 11, 0.5, '2026年中国国际大学生创新大赛', 16, GRAY, False, PP_ALIGN.CENTER)
    page_num(s, 1)

    # ============================================================
    # SLIDE 2: PROJECT OVERVIEW
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, 0.08, SLIDE_H, C_BLUE)
    section_title(s, '项目概述')
    tb(s, 0.8, 1.35, 11.5, 1.8,
       'TCM-TargetMiner — 中药多靶点智能预测与分子对接平台。\n'
       '以"从一味中药发现一个全新药物靶标"为核心理念，整合网络药理学、\n'
       '分子对接、PPI分析和文献挖掘，实现输入中药 → 自动输出候选靶点\n'
       '及 3D 对接结果的全流程闭环。', 15, LIGHT)

    cards = [
        ('痛点', '中药"多成分-多靶点"机制研究难\n现有工具碎片化，流程割裂\n无统一可视化，效率低', C_BLUE),
        ('方案', '输入一味中药→秒出靶点+3D图\n六步全流程在线闭环\n网页端部署，无需VPN', C_GREEN),
        ('优势', '全流程整合（全球首个）\nVina+Fpocket真口袋对接\n独家苗药数据·论文级可视化', C_PURPLE),
    ]
    for i, (ti, de, co) in enumerate(cards):
        x = 0.8 + i * 4.1
        rc(s, x, 3.8, 3.8, 3.0, BG_CARD, co)
        rc(s, x, 3.8, 3.8, 0.06, co)
        tb(s, x + 0.25, 4.1, 1.5, 0.5, ti, 20, co, True)
        tb(s, x + 0.25, 4.8, 3.3, 1.8, de, 13, LIGHT)
    page_num(s, 2)

    # ============================================================
    # SLIDE 3: 6 STEPS (3x2 grid)
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, '核心功能 · 六步全流程')

    steps = [
        ('01 成分筛选', '83味中药 + 8味苗药\nADME自动筛选\nOB≥30%, DL≥0.18', C_BLUE),
        ('02 靶点预测', 'ECFP4分子指纹\n相似度搜索 + 新颖性评分\n排序输出候选靶点', C_GREEN),
        ('03 PPI推断', 'STRING v11数据库\nGuilt-by-Association\n疾病关联智能推断', C_ORANGE),
        ('04 口袋评估', 'Fpocket检测口袋\n体积·疏水性·极性打分\n≥0.5为可成药', C_PURPLE),
        ('05 分子对接', 'AutoDock Vina引擎\n精准对接+结合能计算\n结果自动回填', C_PINK),
        ('06 3D可视化', '论文级3D渲染\n氢键检测(2.5-3.2Å)\n残基自动标注', C_TEAL),
    ]
    for i, (ti, de, co) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = 0.6 + col * 4.2
        y = 1.5 + row * 2.85
        rc(s, x, y, 3.9, 2.5, BG_CARD, co)
        rc(s, x, y, 3.9, 0.06, co)
        # Step number circle
        tb(s, x + 0.2, y + 0.2, 2.5, 0.45, ti, 16, co, True)
        tb(s, x + 0.2, y + 0.8, 3.5, 1.5, de, 13, LIGHT)
    page_num(s, 3)

    # ============================================================
    # SLIDE 4: CASE STUDY - 半枝莲
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, '案例验证 · 半枝莲 (Scutellaria barbata)')

    rc(s, 0.8, 1.4, 5.5, 5.2, BG_CARD, C_BLUE)
    multi_tb(s, 1.0, 1.6, 5.1, 4.8, [
        ('半枝莲', 20, WHITE, True, PP_ALIGN.LEFT),
        ('清热解毒药 · 贵州道地药材', 12, GRAY, False, PP_ALIGN.LEFT),
        ('', 6, LIGHT),
        ('输入成分', 14, C_BLUE, True),
        ('13个化学成分进入分析流程', 13, LIGHT),
        ('', 6, LIGHT),
        ('ADME筛选', 14, C_BLUE, True),
        ('10个满足 OB≥30%, DL≥0.18', 13, LIGHT),
        ('', 6, LIGHT),
        ('靶点预测', 14, C_BLUE, True),
        ('18个候选靶点（含6个全新靶点）', 13, LIGHT),
        ('', 6, LIGHT),
        ('分子对接验证', 14, C_BLUE, True),
        ('Vina对接 → 结合能均＜-6 kcal/mol', 13, C_GREEN, True),
    ])

    rc(s, 6.8, 1.4, 5.7, 5.2, BG_CARD, C_GREEN)
    multi_tb(s, 7.0, 1.6, 5.3, 4.8, [
        ('靶点分析结果', 18, C_GREEN, True),
        ('', 6, LIGHT),
        ('前5关键靶点', 14, WHITE, True),
        ('MAPK3 · MAPK1 · PIK3CA · PTGS2 · AKT1', 14, C_GREEN, True),
        ('', 8, LIGHT),
        ('核心信号通路', 14, WHITE, True),
        ('PI3K-Akt · MAPK · VEGF', 14, LIGHT),
        ('', 8, LIGHT),
        ('推断疾病关联', 14, WHITE, True),
        ('肝细胞癌 · 乳腺癌 · 炎症', 14, LIGHT),
        ('', 8, LIGHT),
        ('功能验证', 14, WHITE, True),
        ('12/18靶点具备可成药口袋', 14, C_TEAL, True),
        ('9/18 对接有效（＜-6 kcal/mol）', 14, C_TEAL, True),
    ])
    page_num(s, 4)

    # ============================================================
    # SLIDE 5: COMPETITION
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, '竞品对比')

    # Header
    headers = ['功能', 'TCMSP', 'SwissTarget', 'HERB', '智鉴药靶']
    hcolors = [BG_CARD, BG_CARD, BG_CARD, BG_CARD, BG_HIGHLIGHT]
    for j, h in enumerate(headers):
        x = 0.8 + j * 2.5
        rc(s, x, 1.5, 2.4, 0.65, hcolors[j], C_BLUE if j < 4 else C_GREEN)
        tb(s, x, 1.55, 2.4, 0.55, h, 15, WHITE if j == 4 else C_BLUE, True, PP_ALIGN.CENTER)

    rows = [
        ['成分查询', '✓', '✗', '✓', '✓+苗药'],
        ['ADME筛选', '✓', '✗', '✓', '✓'],
        ['PPI推断', '✗', '✗', '✗', '✓'],
        ['口袋评估', '✗', '✗', '✗', '✓'],
        ['分子对接', '✗', '✗', '✗', '✓'],
        ['3D可视化', '✗', '✗', '✗', '✓'],
        ['复方/苗药', '✗', '✗', '✗', '✓'],
    ]
    for i, row in enumerate(rows):
        y = 2.25 + i * 0.55
        bg = BG_CARD if i % 2 == 0 else RGBColor(0x10, 0x18, 0x30)
        for j, cell in enumerate(row):
            x = 0.8 + j * 2.5
            cbg = BG_HIGHLIGHT if j == 4 else bg
            rc(s, x, y, 2.4, 0.5, cbg)
            if j == 4:
                clr = C_GREEN if cell == '✓' else (WHITE if cell else LIGHT)
            else:
                clr = C_GREEN if (j > 0 and cell == '✓') else (C_ORANGE if cell == '✗' else LIGHT)
            tb(s, x, y + 0.02, 2.4, 0.46, cell, 13, clr, j == 4, PP_ALIGN.CENTER)

    # 底部总结行
    tb(s, 0.8, 6.3, 12, 0.7,
       '智鉴药靶是唯一覆盖全流程（成分→靶点→对接→可视化）的中药靶点挖掘平台',
       14, C_GREEN, True, PP_ALIGN.CENTER)
    page_num(s, 5)

    # ============================================================
    # SLIDE 6: SWOT
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, 'SWOT 分析')

    swot = [
        ('S', '优势 (Strengths)', '① 全流程整合，功能覆盖度远超同类\n② 论文级3D可视化，直接产出科研成果\n③ 独家苗药数据库，差异化竞争\n④ 开源免费+网页操作，零门槛', C_BLUE),
        ('W', '劣势 (Weaknesses)', '① 平台早期阶段，知名度较低\n② 数据库规模(83味)有待扩充\n③ 团队规模小，人力有限', C_ORANGE),
        ('O', '机会 (Opportunities)', '① "十四五"中医药发展政策支持\n② AI+药物研发市场热潮\n③ 贵州省苗药产业战略机遇\n④ 中医药国际化大趋势', C_GREEN),
        ('T', '威胁 (Threats)', '① 大型生信平台可能进入此领域\n② 开源社区可能出现同类竞品\n③ 数据库持续更新维护的挑战', C_PURPLE),
    ]
    for i, (lb, ti, de, co) in enumerate(swot):
        x = 0.8 + (i % 2) * 6.2
        y = 1.5 + (i // 2) * 2.85
        rc(s, x, y, 5.8, 2.55, BG_CARD, co)
        tb(s, x + 0.25, y + 0.1, 0.5, 0.45, lb, 28, co, True)
        tb(s, x + 0.8, y + 0.15, 4, 0.4, ti, 16, co, True)
        tb(s, x + 0.25, y + 0.6, 5.3, 1.8, de, 13, LIGHT)
    page_num(s, 6)

    # ============================================================
    # SLIDE 7: BUSINESS MODEL
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, '商业模式')

    tiers = [
        ('基础版', '永久免费', '学生/学者免费使用\n全部核心功能\n学术研究用途', C_GREEN),
        ('高校版', '12,800元/年', '中医药院校批量分析\nAPI接口+定制报告\n统一授权管理', C_BLUE),
        ('企业版', '2,980-20,000元/年', '中药企业按规模定价\n专属靶点筛选服务\n研发项目合作', C_ORANGE),
        ('散客', '周18/月30/年280', '独立研究者灵活使用\n按需购买，无需签约\n适合个人课题', C_PURPLE),
    ]
    for i, (ti, price, de, co) in enumerate(tiers):
        x = 0.8 + i * 3.15
        rc(s, x, 1.5, 2.9, 5.3, BG_CARD, co)
        rc(s, x, 1.5, 2.9, 0.06, co)
        tb(s, x + 0.2, 1.75, 2.5, 0.5, ti, 20, co, True, PP_ALIGN.CENTER)
        tb(s, x + 0.2, 2.3, 2.5, 0.5, price, 22, WHITE, True, PP_ALIGN.CENTER)
        ac(s, x + 0.6, 2.9, 1.7, co)
        tb(s, x + 0.2, 3.2, 2.5, 3.0, de, 13, LIGHT, False, PP_ALIGN.CENTER)
    page_num(s, 7)

    # ============================================================
    # SLIDE 8: GLOBAL ROADMAP
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, '全球推广路线')

    # 连接线（时间轴）
    rc(s, 0.8, 3.4, 11.7, 0.04, C_GREEN)

    phases = [
        ('2026-2027', '深耕贵州', '3所高校\n500味药材\n完成软著申请', C_BLUE, 0.8),
        ('2027-2028', '全国拓展', '80+中医药院校\n2000味药材\n英文版上线', C_GREEN, 0.8 + 3.15),
        ('2028-2029', '东亚市场', '日韩汉方市场\n40+高校签约\n15+企业合作', C_ORANGE, 0.8 + 3.15 * 2),
        ('2029-2030', '全球平台', '一带一路沿线推广\n15万注册用户\n多医学体系平台', C_PURPLE, 0.8 + 3.15 * 3),
    ]
    for yr, title, de, co, x in phases:
        # 时间轴节点（圆点用矩形代替）
        rc(s, x + 1.25, 3.28, 0.35, 0.28, co)
        rc(s, x, 1.5, 2.9, 1.6, BG_CARD, co)
        rc(s, x, 1.5, 2.9, 0.05, co)
        tb(s, x + 0.15, 1.7, 2.6, 0.4, yr, 18, co, True, PP_ALIGN.CENTER)
        tb(s, x + 0.15, 2.1, 2.6, 0.35, title, 16, WHITE, True, PP_ALIGN.CENTER)
        rc(s, x, 3.7, 2.9, 2.5, BG_CARD, co)
        tb(s, x + 0.15, 3.9, 2.6, 2.1, de, 13, LIGHT, False, PP_ALIGN.CENTER)

    tb(s, 0.8, 6.5, 11.7, 0.5,
       '目标：2030年建成覆盖全球传统医药体系的智能靶点预测平台',
       14, GRAY, False, PP_ALIGN.CENTER)
    page_num(s, 8)

    # ============================================================
    # SLIDE 9: INNOVATION
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, '核心创新点')

    innov = [
        ('整合创新', '六步在线全流程整合\n全球首个中药靶点挖掘\n一站式服务平台', C_BLUE),
        ('方法学创新', 'Fpocket+Vina真实口袋\n对接，取代传统估算\n结合能的间接方法', C_GREEN),
        ('数据创新', '独家收录8味贵州苗药\n填补民族药物靶点\n研究的空白领域', C_ORANGE),
        ('可视化创新', '论文发表级3D渲染\n氢键自动检测+残基\n标注，结果即插即用', C_PURPLE),
    ]
    for i, (ti, de, co) in enumerate(innov):
        x = 0.8 + i * 3.15
        rc(s, x, 2.0, 2.9, 3.8, BG_CARD, co)
        rc(s, x, 2.0, 2.9, 0.06, co)
        tb(s, x + 0.2, 2.3, 2.5, 0.5, ti, 22, co, True, PP_ALIGN.CENTER)
        rc(s, x + 0.5, 3.0, 1.9, 0.03, co)
        tb(s, x + 0.2, 3.3, 2.5, 2.3, de, 14, LIGHT, False, PP_ALIGN.CENTER)
    page_num(s, 9)

    # ============================================================
    # SLIDE 10: TEAM
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, '团队介绍')

    team = [
        ('项目负责人', '周鑫芸', '平台架构·核心算法\n全栈开发·项目管理', C_BLUE),
        ('技术总监', '', '对接引擎优化\n数据库运维·部署', C_GREEN),
        ('数据科学家', '', '中药数据挖掘\n网络药理学分析', C_TEAL),
        ('市场总监', '', '高校企业合作\n品牌推广·渠道', C_ORANGE),
        ('产品经理', '', '需求调研\n迭代规划·测试', C_PURPLE),
        ('运营财务', '', '平台运营\n预算编制·行政', C_PINK),
    ]
    for i, (role, name, desc, co) in enumerate(team):
        col = i % 3
        row = i // 3
        x = 0.8 + col * 4.1
        y = 1.5 + row * 2.85
        rc(s, x, y, 3.8, 2.45, BG_CARD, co)
        rc(s, x, y, 3.8, 0.06, co)
        tb(s, x + 0.2, y + 0.25, 3.4, 0.4, role, 18, co, True, PP_ALIGN.CENTER)
        if name:
            tb(s, x + 0.2, y + 0.7, 3.4, 0.35, name, 15, WHITE, True, PP_ALIGN.CENTER)
        tb(s, x + 0.2, y + (1.15 if name else 0.75), 3.4, 1.3, desc, 13, LIGHT, False, PP_ALIGN.CENTER)
    page_num(s, 10)

    # ============================================================
    # SLIDE 11: FINANCE
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    section_title(s, '财务预测')

    tb(s, 0.8, 1.3, 11.5, 0.7,
       '每年盈利，不烧钱。零边际成本模型保障高利润率。', 18, LIGHT)

    # 营收增长可视化条
    rev_data = [
        ('2026', '3.84万', 0.08, C_BLUE),
        ('2027', '19.2万', 0.38, C_GREEN),
        ('2028', '51.2万+', 1.0, C_TEAL),
    ]
    y_start = 2.3
    for i, (yr, rev, ratio, co) in enumerate(rev_data):
        y = y_start + i * 0.75
        bar_w = 6.5 * ratio
        tb(s, 0.8, y, 1.0, 0.5, yr, 16, WHITE, True)
        rc(s, 1.8, y + 0.08, 6.5, 0.35, RGBColor(0x1A, 0x28, 0x48))
        rc(s, 1.8, y + 0.08, max(bar_w, 0.5), 0.35, co)
        tb(s, 1.8 + max(bar_w, 0.5) + 0.15, y, 2, 0.5, rev, 16, co, True)

    # 2029-2030 预计
    y = y_start + 3 * 0.75
    tb(s, 0.8, y, 1.0, 0.5, '2029-30', 16, WHITE, True)
    tb(s, 1.8, y + 0.08, 6.5, 0.5, '15万用户 + 企业版（预计营收突破200万）', 15, C_GREEN, True)

    # 成本结构
    rc(s, 0.8, 5.2, 5.5, 1.8, BG_CARD, C_GREEN)
    multi_tb(s, 1.0, 5.4, 5.1, 1.6, [
        ('成本结构', 16, C_GREEN, True),
        ('开源技术栈 → 零软件授权费', 13, LIGHT),
        ('学校计算资源 → 零硬件成本', 13, LIGHT),
        ('云端部署 → 按需付费，边际成本极低', 13, LIGHT),
        ('综合利润率 > 70%', 14, WHITE, True),
    ])

    rc(s, 6.8, 5.2, 5.7, 1.8, BG_CARD, C_BLUE)
    multi_tb(s, 7.0, 5.4, 5.3, 1.6, [
        ('里程碑', 16, C_BLUE, True),
        ('2026: 500味药材 · 3所高校 · 软著', 13, LIGHT),
        ('2027: 2000味 · MD模拟 · 英文版', 13, LIGHT),
        ('2028: 日韩版 · 12人团队 · 40+高校', 13, LIGHT),
        ('2030: 全球多医学体系平台', 13, WHITE, True),
    ])
    page_num(s, 11)

    # ============================================================
    # SLIDE 12: THANK YOU
    # ============================================================
    s = prs.slides.add_slide(sl)
    dark_bg(s)
    rc(s, 0, 0, SLIDE_W, 0.08, C_BLUE)
    rc(s, 0, SLIDE_H - 0.08, SLIDE_W, 0.08, C_GREEN)
    tb(s, 1, 2.2, 11, 1.5, '智鉴药靶', 52, WHITE, True, PP_ALIGN.CENTER)
    ac(s, 5.2, 3.7, 2.933, C_GREEN)
    tb(s, 1, 3.9, 11, 0.7, '感谢聆听 · 敬请指导', 26, C_BLUE, False, PP_ALIGN.CENTER)
    tb(s, 1, 4.9, 11, 0.5, '周鑫芸 | 贵州中医药大学 药学院', 18, GRAY, False, PP_ALIGN.CENTER)
    page_num(s, 12)

    return prs


def generate_bytes():
    buf = io.BytesIO()
    build_prs().save(buf)
    buf.seek(0)
    return buf.read()


def save_to(path):
    build_prs().save(path)
    return path
